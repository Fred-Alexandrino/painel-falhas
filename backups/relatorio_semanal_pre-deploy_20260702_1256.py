# -*- coding: utf-8 -*-
"""
relatorio_semanal.py
─────────────────────────────────────────────────────────────────────────
Geração automática do "Relatório Semanal" (.pptx) no modelo Grid Co.,
a partir das ocorrências já registradas na planilha do Painel de Falhas.

Não usa nenhuma API de IA — o texto narrativo é montado por regras
determinísticas em cima dos campos: Falha, Causa, Ação e Histórico.

Como plugar no app.py:
    from relatorio_semanal import gerar_relatorio_pptx
    (ver bloco de rota Flask no final deste arquivo)
"""

import os
import re
import copy
import unicodedata
from io import BytesIO
from datetime import datetime

from pptx import Presentation
from pptx.oxml.ns import qn

# ── Configuração ───────────────────────────────────────────────────────────

# Caminho do modelo Grid Co. (subir este arquivo no repo, ex: templates/modelo_relatorio_semanal.pptx)
TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "templates", "modelo_relatorio_semanal.pptx")

# Índices das colunas da planilha (0-based, iguais ao app.py / CAMPO_COL)
COL_ID, COL_CLIENTE, COL_USINA, COL_EQUIP = 0, 1, 2, 3
COL_FALHA, COL_CAUSA, COL_IMPACTADOS, COL_ACAO = 4, 5, 6, 7
COL_STATUS, COL_TICKET, COL_OS, COL_HISTORICO = 8, 9, 10, 11
COL_DATA_ABERTURA = 12
COL_DATA_FECHAMENTO = 20  # coluna U (21ª, 1-based) = índice 20 (0-based)

STATUS_CONCLUIDO = ["concluído", "concluido", "resolvido", "fechado", "resolved", "closed"]

# Categorias "bonitas" para exibir no relatório (padrões regex, avaliados em ordem —
# o primeiro que bater define a categoria). Casa tanto nomes por extenso ("Inversor")
# quanto códigos curtos usados em campo ("Inv-18", "NVR", "Chave Seccionadora Skid 1").
PADROES_CATEGORIA = [
    (r"\binv[\s\-]*\d+|inversor", "INVERSORES"),
    (r"tracker|\btcu\b", "TRACKERS / TCU"),
    (r"fusive|fusíve", "FUSÍVEIS"),
    (r"transformador", "TRANSFORMADORES"),
    (r"switchgear|chave seccionadora|disjuntor|religador", "SWITCHGEAR"),
    (r"\bstring\b|modulo|módulo|otimizador", "STRINGS / MÓDULOS"),
    (r"usina desligad|usina offline|desligamento|ufv desligad|usina parad", "DESLIGAMENTOS"),
    (r"comunica|scada|cctv|\bnvr\b|camera|câmera|speed dome|igate", "COMUNICAÇÕES / SCADA / CCTV"),
    (r"sensor|piranometro|piranômetro|solarimetric|estacao solarimetrica|estação solarimétrica", "SENSORES"),
    (r"emergenc", "EMERGÊNCIAS"),
    (r"operac|opera[çc][aã]o", "OPERAÇÕES"),
    (r"termografia|termograf", "TERMOGRAFIA"),
    (r"restart|religamento", "RESTART / RELIGAMENTO"),
    (r"ronda", "RONDAS SEMANAIS"),
    (r"exaustor", "EXAUSTORES"),
    (r"vegeta", "CONTROLE DE VEGETAÇÃO"),
]


# ── Utilidades de texto/data ────────────────────────────────────────────────

def _norm(s):
    s = unicodedata.normalize("NFKD", (s or "").lower())
    s = s.encode("ascii", "ignore").decode("ascii")
    return s.strip()


def _rotulo_categoria(equip_bruto):
    n = _norm(equip_bruto)
    for padrao, rotulo in PADROES_CATEGORIA:
        if re.search(padrao, n):
            return rotulo
    return (equip_bruto or "OUTRAS OCORRÊNCIAS").strip().upper()


def _parse_data(txt):
    """Aceita 'dd/mm/aaaa[ hh:mm[:ss]]' -> datetime, ou None."""
    if not txt:
        return None
    txt = txt.strip()
    for fmt in ("%d/%m/%Y %H:%M:%S", "%d/%m/%Y %H:%M", "%d/%m/%Y"):
        try:
            return datetime.strptime(txt, fmt)
        except ValueError:
            continue
    return None


def _is_concluido(status):
    s = (status or "").lower()
    return any(x in s for x in STATUS_CONCLUIDO)


# ── Coleta e agrupamento dos dados ──────────────────────────────────────────

def coletar_ocorrencias_semana(todos_valores, cliente, data_inicio, data_fim):
    """
    todos_valores: retorno de ws.get_all_values() (linha 0 = cabeçalho)
    cliente: nome do cliente (comparação por 'contains', case-insensitive)
    data_inicio / data_fim: datetime (00:00 do dia inicial até 23:59 do dia final)

    Retorna dict:
        {
          "categoria": {
              "concluidas": [linha, linha, ...],
              "abertas":    [linha, linha, ...],
          }, ...
        }
    """
    cliente_norm = _norm(cliente)
    grupos = {}

    for row in todos_valores[1:]:
        if len(row) <= COL_DATA_FECHAMENTO:
            row = row + [""] * (COL_DATA_FECHAMENTO + 1 - len(row))

        if cliente_norm not in _norm(row[COL_CLIENTE]):
            continue

        dt_abertura = _parse_data(row[COL_DATA_ABERTURA])
        dt_fecham = _parse_data(row[COL_DATA_FECHAMENTO])
        concluida = _is_concluido(row[COL_STATUS])

        fechou_na_semana = dt_fecham and data_inicio <= dt_fecham <= data_fim
        abriu_na_semana = dt_abertura and data_inicio <= dt_abertura <= data_fim
        em_aberto_backlog = not concluida

        if not (fechou_na_semana or abriu_na_semana or em_aberto_backlog):
            continue

        categoria = _rotulo_categoria(row[COL_EQUIP])
        grupos.setdefault(categoria, {"concluidas": [], "abertas": []})

        if fechou_na_semana:
            grupos[categoria]["concluidas"].append(row)
        elif em_aberto_backlog:
            grupos[categoria]["abertas"].append(row)

    # remove categorias que ficaram vazias
    return {k: v for k, v in grupos.items() if v["concluidas"] or v["abertas"]}


# ── Geração da narrativa (offline, sem IA) ─────────────────────────────────

def _frase_ocorrencia(row, concluida):
    usina = row[COL_USINA].strip()
    falha = row[COL_FALHA].strip()
    causa = row[COL_CAUSA].strip()
    acao = row[COL_ACAO].strip()
    historico = row[COL_HISTORICO].strip()
    dt_abertura = _parse_data(row[COL_DATA_ABERTURA])
    dt_fecham = _parse_data(row[COL_DATA_FECHAMENTO])

    partes = []

    abertura_txt = f" em {dt_abertura.strftime('%d/%m')}" if dt_abertura else ""
    if falha:
        partes.append(f"Foi registrada ocorrência na usina {usina}{abertura_txt}: {falha}.")
    else:
        partes.append(f"Foi registrada ocorrência na usina {usina}{abertura_txt}.")

    if causa:
        partes.append(f"A causa identificada foi {causa[0].lower()}{causa[1:]}.")

    if acao:
        partes.append(f"A equipe Grid Co. atuou com: {acao[0].lower()}{acao[1:]}.")

    if concluida:
        fecham_txt = f" em {dt_fecham.strftime('%d/%m')}" if dt_fecham else ""
        partes.append(f"O caso foi normalizado{fecham_txt}, com a operação restabelecida.")
    else:
        partes.append("O caso segue em acompanhamento pela equipe até a normalização completa.")

    return " ".join(partes)


def gerar_texto_categoria(categoria, dados):
    """Monta o parágrafo narrativo de uma categoria, no tom Grid Co."""
    frases = []
    concluidas = dados["concluidas"]
    abertas = dados["abertas"]

    if concluidas:
        if len(concluidas) == 1:
            frases.append(_frase_ocorrencia(concluidas[0], concluida=True))
        else:
            frases.append(
                f"Na semana, foram normalizadas {len(concluidas)} ocorrências relacionadas a "
                f"{categoria.lower()}, com destaque para os seguintes atendimentos:"
            )
            for r in concluidas:
                frases.append("• " + _frase_ocorrencia(r, concluida=True))

    if abertas:
        if concluidas:
            frases.append("Seguem em acompanhamento, ainda em aberto:")
        if len(abertas) == 1:
            frases.append(_frase_ocorrencia(abertas[0], concluida=False))
        else:
            for r in abertas:
                frases.append("• " + _frase_ocorrencia(r, concluida=False))

    return "\n".join(frases)


# ── Clonagem de slides (mantém 100% do estilo do modelo Grid Co.) ──────────

def _duplicate_slide(prs, index):
    """Duplica um slide existente do modelo (mesmo layout, cores, fontes, formas)."""
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)

    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    id_map = {}
    for rId, rel in source.part.rels.items():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        new_rid = dest.part.rels._add_relationship(rel.reltype, rel._target, rel.is_external)
        id_map[rId] = new_rid

    for shp in source.shapes:
        newel = copy.deepcopy(shp._element)
        for el in newel.iter():
            for attr in ("embed", "link", "id"):
                full = qn(f"r:{attr}")
                val = el.get(full)
                if val and val in id_map:
                    el.set(full, id_map[val])
        dest.shapes._spTree.append(newel)
    return dest


def _set_text_preservando_estilo(shape, novo_texto):
    """Substitui o texto de um text box mantendo a formatação do 1º run."""
    tf = shape.text_frame
    primeiro_par = tf.paragraphs[0]
    if primeiro_par.runs:
        modelo_run = primeiro_par.runs[0]
        font = modelo_run.font
        cor_rgb = None
        try:
            if font.color and font.color.type is not None:
                cor_rgb = font.color.rgb  # levanta AttributeError se for cor de tema (schemeClr)
        except AttributeError:
            cor_rgb = None
        fmt = (font.name, font.size, font.bold, font.italic, cor_rgb)
    else:
        fmt = None

    tf.clear()
    linhas = novo_texto.split("\n")
    for i, linha in enumerate(linhas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = linha
        if fmt:
            if fmt[0]:
                run.font.name = fmt[0]
            if fmt[1]:
                run.font.size = fmt[1]
            if fmt[2] is not None:
                run.font.bold = fmt[2]
            if fmt[3] is not None:
                run.font.italic = fmt[3]
            if fmt[4]:
                run.font.color.rgb = fmt[4]


def _remover_fotos(slide):
    """Remove os placeholders de foto (o gerador não tem fotos de campo)."""
    for shp in list(slide.shapes):
        if shp.shape_type == 13 and (shp.name or "").lower().startswith("imagem"):
            shp._element.getparent().remove(shp._element)


def _find_shape(slide, contem_texto):
    for shp in slide.shapes:
        if shp.has_text_frame and contem_texto.lower() in shp.text_frame.text.lower():
            return shp
    return None


# ── Montagem do PPTX final ──────────────────────────────────────────────────

def agrupar_em_paginas(grupos, orcamento_chars=1400):
    """
    Agrupa categorias em páginas por orçamento de caracteres, para não desperdiçar
    slide com pouco conteúdo nem lotar demais uma página. Retorna lista de páginas,
    cada página = lista de tuplas (categoria, texto).
    """
    blocos = [(cat, gerar_texto_categoria(cat, dados)) for cat, dados in grupos.items()]
    paginas, atual, atual_len = [], [], 0
    for cat, texto in blocos:
        tam = len(texto) + len(cat) + 10
        if atual and (atual_len + tam > orcamento_chars):
            paginas.append(atual)
            atual, atual_len = [], 0
        atual.append((cat, texto))
        atual_len += tam
    if atual:
        paginas.append(atual)
    return paginas


def gerar_relatorio_pptx(cliente, usinas_label, semana_label, grupos):
    """
    grupos: dict categoria -> {"concluidas": [...], "abertas": [...]}
    Retorna BytesIO() pronto para download.
    """
    prs = Presentation(TEMPLATE_PATH)

    # --- Slide 1: Capa (clona o slide 0 do modelo e troca o texto) --------
    capa = _duplicate_slide(prs, 0)
    shp_titulo = _find_shape(capa, "O&M")
    if shp_titulo:
        _set_text_preservando_estilo(
            shp_titulo, f"O&M – {cliente}\n{usinas_label}\n{semana_label}"
        )

    # --- Slide 2: Ata da reunião (somente os tópicos/categorias da semana) -
    ata = _duplicate_slide(prs, 1)
    shp_lista = _find_shape(ata, "Desligamentos")  # caixa de texto com a lista de tópicos
    if shp_lista:
        _set_text_preservando_estilo(shp_lista, "\n".join(grupos.keys()))

    # --- Uma página por grupo de categorias (empacotadas por orçamento) ----
    for pagina in agrupar_em_paginas(grupos):
        novo = _duplicate_slide(prs, 2)
        _remover_fotos(novo)

        shp_categoria = _find_shape(novo, "DESLIGAMENTOS")
        shp_corpo = _find_shape(novo, "Foram registradas")

        if len(pagina) == 1:
            categoria, texto = pagina[0]
            if shp_categoria:
                _set_text_preservando_estilo(shp_categoria, categoria)
            if shp_corpo:
                _set_text_preservando_estilo(shp_corpo, texto)
        else:
            if shp_categoria:
                _set_text_preservando_estilo(shp_categoria, "OCORRÊNCIAS DA SEMANA")
            if shp_corpo:
                blocos_txt = [f"{cat}\n{texto}" for cat, texto in pagina]
                _set_text_preservando_estilo(shp_corpo, "\n\n".join(blocos_txt))

    # --- Reordena o deck: capa nova -> ata nova -> categorias -> contato --
    # sldIdLst atual: [0]capa-orig [1]ata-orig [2..5]conteudo-orig [6]contato-orig
    #                 [7]capa-clone [8]ata-clone [9:]categorias-clone (nesta ordem de criação)
    xml_slides = prs.slides._sldIdLst
    todos_els = list(xml_slides)
    contato_el = todos_els[6]
    capa_el = todos_els[7]
    ata_el = todos_els[8]
    categorias_els = todos_els[9:]

    for e in todos_els:
        xml_slides.remove(e)
    for e in [capa_el, ata_el] + categorias_els + [contato_el]:
        xml_slides.append(e)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── Rota Flask (adicionar ao app.py) ────────────────────────────────────────
"""
from flask import request, send_file
from relatorio_semanal import coletar_ocorrencias_semana, gerar_relatorio_pptx
from datetime import datetime

@app.route("/gerar-relatorio-semanal", methods=["POST", "OPTIONS"])
def gerar_relatorio_semanal_route():
    if request.method == "OPTIONS":
        return ("", 204)
    try:
        body = request.get_json(force=True) or {}
        cliente = str(body.get("cliente", "")).strip()
        data_inicio = datetime.strptime(body["dataInicio"], "%Y-%m-%d")
        data_fim = datetime.strptime(body["dataFim"], "%Y-%m-%d").replace(
            hour=23, minute=59, second=59
        )
        if not cliente:
            return jsonify({"ok": False, "error": "cliente é obrigatório"}), 400

        todos = carregar_planilha(get_sheet())
        grupos = coletar_ocorrencias_semana(todos, cliente, data_inicio, data_fim)
        if not grupos:
            return jsonify({"ok": False, "error": "Nenhuma ocorrência encontrada no período"}), 404

        semana_label = f"{data_inicio.strftime('%d/%m')} a {data_fim.strftime('%d/%m/%Y')}"
        buf = gerar_relatorio_pptx(cliente, cliente, semana_label, grupos)

        nome_arquivo = f"Relatorio_{cliente}_{data_inicio.strftime('%Y%m%d')}.pptx".replace(" ", "_")
        return send_file(
            buf,
            as_attachment=True,
            download_name=nome_arquivo,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except Exception as e:
        log.error(f"[Relatorio Semanal] Erro: {e}")
        return jsonify({"ok": False, "error": str(e)}), 500
"""
