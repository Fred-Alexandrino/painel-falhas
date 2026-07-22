# -*- coding: utf-8 -*-
"""
relatorio_semanal.py
─────────────────────────────────────────────────────────────────────────
Geração automática do "Relatório Semanal" (.pptx) no modelo Grid Co.,
a partir das ocorrências (Painel de Falhas) E das OSs/atividades (Painel
de Atividades) já registradas nas planilhas.

Não usa nenhuma API de IA — o texto é montado por regras determinísticas
em cima dos campos de cada fonte (resumido, sem despejar o histórico
bruto de ações).
"""

import os
import re
import copy
import unicodedata
from io import BytesIO
from datetime import datetime

from pptx import Presentation
from pptx.oxml.ns import qn

# ── Configuração ───────────────────────────────────────────────────────────

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "templates", "modelo_relatorio_semanal.pptx")

# Índices das colunas da planilha principal — Painel de Falhas (0-based,
# iguais ao app.py / CAMPO_COL)
COL_ID, COL_CLIENTE, COL_USINA, COL_EQUIP = 0, 1, 2, 3
COL_FALHA, COL_CAUSA, COL_IMPACTADOS, COL_ACAO = 4, 5, 6, 7
COL_STATUS, COL_TICKET, COL_OS, COL_HISTORICO = 8, 9, 10, 11
COL_DATA_ABERTURA = 12
COL_DATA_FECHAMENTO = 20  # coluna U (21ª, 1-based) = índice 20 (0-based)

# Índices das colunas da aba Painel de Atividades (0-based, iguais ao
# ATIV_CAMPO_COL do app.py, que é 1-based -> aqui -1)
ATIV_COL_ID, ATIV_COL_CLIENTE, ATIV_COL_USINA, ATIV_COL_EQUIP = 0, 1, 2, 3
ATIV_COL_DESCRICAO, ATIV_COL_RESPONSAVEL, ATIV_COL_PRAZO, ATIV_COL_PRIORIDADE = 4, 5, 6, 7
ATIV_COL_STATUS, ATIV_COL_DATA_CRIACAO, ATIV_COL_DATA_CONCLUSAO, ATIV_COL_HISTORICO = 8, 9, 10, 11
ATIV_COL_EDITOR, ATIV_COL_NUMEROOS, ATIV_COL_STATUSOS = 12, 13, 14

STATUS_CONCLUIDO = ["concluído", "concluido", "resolvido", "fechado", "resolved", "closed"]

# CORRIGIDO 22/07/2026: uma atividade/OS só conta como "realizada" no
# RELATÓRIO SEMANAL (cliente-facing) quando o statusOS chega em
# "Finalizada" — igual à aba "Concluídas" do próprio Painel de Atividades
# (statusOS=="Finalizada"; "Em Revisão" é a aba separada "Em Verificação").
# A versão anterior também tratava "Em Revisão" como concluída aqui, o que
# fazia o relatório mostrar como pronta uma atividade ainda em validação
# (a fallback de "abriu E já está 'concluída'" pegava OSs criadas na
# semana que só tinham chegado em Em Revisão) — Fred precisou corrigir
# manualmente o relatório da Alves Lima (Semana 30) por causa disso.
# Esta definição é local a este módulo — não afeta as abas/KPIs do
# dashboard em app.py (que usam sua própria lógica de Em Verificação).
STATUS_OS_CONCLUIDO = ["finalizada"]

# Falhas já convertidas em Atividade não devem aparecer de novo como
# ocorrência "em aberto" no relatório — o trabalho já é reportado (com o
# número de OS) através de coletar_atividades_semana(). Sem esta exclusão,
# como "Convertida em Atividade" nunca bate em STATUS_CONCLUIDO, a
# ocorrência cai perpetuamente em "aberta"/backlog e aparece em toda
# semana seguinte, duplicando a mesma informação. Corrigido 22/07/2026.
STATUS_FALHA_EXCLUIR_RELATORIO = ["convertida em atividade", "convertido em atividade"]

CAT_COMUNICACOES = "COMUNICAÇÕES / SCADA / CCTV"
CAT_DESLIGAMENTOS = "DESLIGAMENTOS"

# Pauta padrão da "Ata da reunião" (slide 2) — combinado com Fred em
# 22/07/2026: para todo cliente, EXCETO RenoGrid, a pauta é sempre esta
# lista fixa, independente das categorias que de fato tiveram ocorrência
# na semana. RenoGrid mantém a pauta dinâmica (montada a partir do que
# realmente aparece no relatório daquela semana).
PAUTA_PADRAO_ATA = ["ATIVIDADES SEMANAIS", CAT_DESLIGAMENTOS, "PCM"]
CLIENTES_PAUTA_DINAMICA = ["renogrid"]

# Categorias "bonitas" para exibir no relatório (padrões regex, avaliados em ordem —
# o primeiro que bater define a categoria). Casa tanto nomes por extenso ("Inversor")
# quanto códigos curtos usados em campo ("Inv-18", "NVR", "Chave Seccionadora Skid 1").
# A ordem desta lista também define a ordem de exibição das categorias no relatório
# (exceto Comunicações e Desligamentos, que sempre recebem página própria — ver
# gerar_relatorio_pptx).
PADROES_CATEGORIA = [
    (r"\binv[\s\-]*r?\.?\d+|inversor", "INVERSORES"),
    (r"tracker|\btcu\b", "TRACKERS / TCU"),
    (r"fusive|fusíve", "FUSÍVEIS"),
    (r"transformador", "TRANSFORMADORES"),
    (r"switchgear|chave seccionadora|disjuntor|religador", "SWITCHGEAR"),
    (r"\bskid", "SKIDS"),
    (r"preventiv", "MANUTENÇÃO PREVENTIVA"),
    (r"\bstring\b|modulo|módulo|otimizador", "STRINGS / MÓDULOS"),
    (r"usina desligad|usina offline|desligamento|ufv desligad|usina parad", CAT_DESLIGAMENTOS),
    (r"comunica|scada|cctv|cftv|\bnvr\b|camera|câmera|speed dome|igate", CAT_COMUNICACOES),
    (r"sensor|piranometro|piranômetro|solarimetric|estacao solarimetrica|estação solarimétrica", "SENSORES"),
    (r"emergenc", "EMERGÊNCIAS"),
    (r"operac|opera[çc][aã]o", "OPERAÇÕES"),
    (r"termografia|termograf", "TERMOGRAFIA"),
    (r"restart|religamento", "RESTART / RELIGAMENTO"),
    (r"ronda", "RONDAS SEMANAIS"),
    (r"exaustor", "EXAUSTORES"),
    (r"vegeta", "CONTROLE DE VEGETAÇÃO"),
]

# Rótulo genérico pra tudo que não bate em nenhum padrão acima (formulários,
# spare parts, reclamações de concessionária, códigos Fracttal sem nenhuma
# palavra reconhecível etc.) — nunca usa o texto bruto do equipamento como
# categoria própria, pra não fragmentar o relatório em uma mini-seção por
# código.
CAT_OUTRAS = "OUTRAS OCORRÊNCIAS"

# Ordem de exibição das categorias "genéricas" (tudo que não é Comunicações
# nem Desligamentos, que são tratadas à parte). Categorias fora desta lista
# (rótulos livres vindos direto do campo Equipamento) vão para o final, em
# ordem alfabética.
ORDEM_CATEGORIAS_GERAIS = [rotulo for _, rotulo in PADROES_CATEGORIA
                           if rotulo not in (CAT_COMUNICACOES, CAT_DESLIGAMENTOS)] + [CAT_OUTRAS]

# Zeladoria: colunas fixas da tabela do relatório
ZELADORIA_COLUNAS = ["Usina", "Roçagem", "Poda Química", "Limpeza dos Módulos"]
# Mapeia cada coluna da tabela para o "nome" de grupo usado em coletar_zeladoria()
ZELADORIA_MAPA_COLUNAS = [("Roçagem", "Roçada"), ("Poda Química", "Poda Química"),
                          ("Limpeza dos Módulos", "Lavagem dos Módulos")]
ZELADORIA_MAX_USINAS_POR_PAGINA = 8  # evita estourar a tabela pra fora do slide


# ── Utilidades de texto/data ────────────────────────────────────────────────

def _norm(s):
    s = unicodedata.normalize("NFKD", (s or "").lower())
    s = s.encode("ascii", "ignore").decode("ascii")
    return s.strip()


def _rotulo_categoria(equip_bruto):
    n = _norm(equip_bruto)
    for padrao, rotulo in PADROES_CATEGORIA:
        if re.search(padrao, n):
            return rotulo
    return CAT_OUTRAS


def _parse_data(txt):
    """Aceita 'dd/mm/aaaa[ hh:mm[:ss]]' -> datetime, ou None."""
    if not txt:
        return None
    txt = txt.strip()
    for fmt in ("%d/%m/%Y %H:%M:%S", "%d/%m/%Y %H:%M", "%d/%m/%Y"):
        try:
            return datetime.strptime(txt, fmt)
        except ValueError:
            continue
    return None


_RE_DATA_HISTORICO = re.compile(r"(\d{2}/\d{2}/\d{4})\s+(\d{2}:\d{2}(?::\d{2})?)")


def _ultima_data_historico(historico_texto):
    """
    Extrai a data/hora da ÚLTIMA linha com timestamp no Histórico da OS
    (formato "dd/mm/aaaa HH:MM - ..."). Corrigido 22/07/2026: a Data de
    Conclusão gravada pelo nosso sistema reflete quando a AUDITORIA
    detectou o "Finalizada" na Fracttal — não quando isso realmente
    aconteceu. Uma OS com trabalho concluído em 10/07 mas só confirmada
    pela auditoria dias depois (dentro da janela do relatório seguinte)
    aparecia como "concluída nesta semana" mesmo sendo antiga (caso
    relatado pelo Fred: OS 9173). O Histórico, gravado a cada mudança real
    de status na Fracttal, é a fonte mais confiável da data verdadeira —
    por isso tem prioridade sobre a Data de Conclusão ao decidir se algo
    pertence ao período do relatório.
    """
    if not historico_texto:
        return None
    melhor = None
    for linha in historico_texto.splitlines():
        m = _RE_DATA_HISTORICO.search(linha)
        if not m:
            continue
        dt = _parse_data(f"{m.group(1)} {m.group(2)}")
        if dt and (melhor is None or dt > melhor):
            melhor = dt
    return melhor


def _fmt_data_hora(dt):
    """datetime -> 'dd/mm/aaaa às HH:MM'."""
    return dt.strftime("%d/%m/%Y às %H:%M")


def _is_concluido(status):
    s = (status or "").lower()
    return any(x in s for x in STATUS_CONCLUIDO)


def _atividade_concluida(status_os):
    return (status_os or "").strip().lower() in STATUS_OS_CONCLUIDO


def _status_exibicao(status_raw):
    """Traduz o status bruto (statusOS do Fracttal / status de Falha) para
    o rótulo mostrado no relatório do cliente — evita imprimir jargão
    interno como 'Finalizada'/'Em Revisão' sem contexto. Corrigido
    22/07/2026 (ver ajuste manual do Fred no relatório Alves Lima Semana 30)."""
    s = (status_raw or "").strip().lower()
    if s == "finalizada":
        return "Concluída"
    if s in ("em revisão", "em revisao"):
        return "Em verificação"
    if s in STATUS_CONCLUIDO:
        return "Concluído"
    return (status_raw or "Em aberto").strip()


def _ticket_valido(ticket):
    """Só considera 'chamado real' se houver um número/identificador de verdade."""
    t = (ticket or "").strip().lower()
    return t not in ("", "n/a", "na", "-", "--", "nao", "não", "sem chamado", "sem ticket", "s/n")


def _ordenar_categorias(chaves):
    """Ordena rótulos de categoria pela ordem definida em ORDEM_CATEGORIAS_GERAIS;
    o que não estiver na lista vai para o final, em ordem alfabética."""
    def chave_ordenacao(rotulo):
        if rotulo in ORDEM_CATEGORIAS_GERAIS:
            return (0, ORDEM_CATEGORIAS_GERAIS.index(rotulo))
        return (1, rotulo)
    return sorted(chaves, key=chave_ordenacao)


# ── Coleta e agrupamento — Painel de Falhas (ocorrências) ──────────────────
#
# Toda ocorrência/atividade, seja qual for a origem, é normalizada para o
# mesmo formato de dict antes de entrar em "grupos":
#   {"usina", "equipamento", "falha", "status", "dt_abertura",
#    "dt_fechamento", "concluida", "origem"}
# Isso permite juntar Painel de Falhas + Painel de Atividades num único
# relatório sem duplicar as funções de geração de texto/slide.

def coletar_ocorrencias_semana(todos_valores, cliente, data_inicio, data_fim):
    """
    todos_valores: retorno de ws.get_all_values() do Painel de Falhas (linha 0 = cabeçalho)
    cliente: nome do cliente (comparação por 'contains', case-insensitive)
    data_inicio / data_fim: datetime (00:00 do dia inicial até 23:59 do dia final)

    CORRIGIDO 22/07/2026: só entra o que abriu OU fechou dentro do período
    selecionado (mesma correção aplicada em coletar_atividades_semana) —
    uma ocorrência aberta há muito tempo não deve reaparecer em todo
    relatório seguinte só por continuar sem status de concluído.

    Retorna dict: {"categoria": {"concluidas": [dict,...], "abertas": [dict,...]}, ...}
    """
    cliente_norm = _norm(cliente)
    grupos = {}

    for row in todos_valores[1:]:
        if len(row) <= COL_DATA_FECHAMENTO:
            row = row + [""] * (COL_DATA_FECHAMENTO + 1 - len(row))

        if cliente_norm not in _norm(row[COL_CLIENTE]):
            continue

        if _norm(row[COL_STATUS]) in STATUS_FALHA_EXCLUIR_RELATORIO:
            # já reportada como Atividade (com número de OS) — evita duplicar
            continue

        dt_abertura = _parse_data(row[COL_DATA_ABERTURA])
        concluida = _is_concluido(row[COL_STATUS])
        if concluida:
            # mesmo raciocínio de coletar_atividades_semana: prioriza a
            # última data real do Histórico sobre a Data de Fechamento
            # (que pode ter sido gravada só quando alguém confirmou o
            # fechamento no painel, não quando o problema foi resolvido
            # de fato).
            dt_fecham = _ultima_data_historico(row[COL_HISTORICO]) or _parse_data(row[COL_DATA_FECHAMENTO])
        else:
            dt_fecham = None

        fechou_na_semana = dt_fecham and data_inicio <= dt_fecham <= data_fim
        abriu_na_semana = dt_abertura and data_inicio <= dt_abertura <= data_fim

        if not (fechou_na_semana or abriu_na_semana):
            continue

        texto_categoria = f"{row[COL_EQUIP]} {row[COL_FALHA]}"
        categoria = _rotulo_categoria(texto_categoria)
        d = {
            "usina": row[COL_USINA].strip() or "Usina não informada",
            "equipamento": row[COL_EQUIP].strip() or "Equipamento não informado",
            "numero_os": row[COL_OS].strip() or None,
            "falha": row[COL_FALHA].strip() or "Falha não especificada",
            "status": row[COL_STATUS].strip() or "Em aberto",
            "dt_abertura": dt_abertura,
            "dt_fechamento": dt_fecham,
            "concluida": concluida,
            "origem": "falha",
        }
        grupos.setdefault(categoria, {"concluidas": [], "abertas": []})

        if fechou_na_semana or (concluida and abriu_na_semana):
            # fechou dentro da semana, OU abriu e já está concluída mas sem
            # Data de Fechamento gravada corretamente (mesmo fallback usado
            # em coletar_atividades_semana).
            grupos[categoria]["concluidas"].append(d)
        else:
            # abriu dentro da semana e ainda está em aberto — mostra como
            # aberta só nesta janela (não perpetuamente).
            grupos[categoria]["abertas"].append(d)

    return {k: v for k, v in grupos.items() if v["concluidas"] or v["abertas"]}


# ── Coleta e agrupamento — Painel de Atividades (OSs) ───────────────────────

def coletar_atividades_semana(todos_valores, cliente, data_inicio, data_fim):
    """
    todos_valores: retorno de ws.get_all_values() do Painel de Atividades
    (linha 0 = cabeçalho). A data de abertura é a Data de Criação da
    OS/atividade e a de fechamento é a Data de Conclusão; "concluída" é
    definido pelo statusOS (só "Finalizada" — ver STATUS_OS_CONCLUIDO).

    CORRIGIDO 22/07/2026: o relatório semanal só deve trazer o que
    realmente aconteceu (abriu OU fechou) DENTRO do período selecionado —
    uma atividade ainda aberta (ex.: presa em "Em Revisão" há semanas,
    criada bem antes do período) não deve reaparecer em todo relatório
    seguinte só por continuar sem "Finalizada". O caso relatado pelo Fred:
    OS 9173 (ABC Morada Nova), criada em 08/07 e parada em "Em Revisão" a
    100% desde então, aparecia num relatório de 16/07 a 22/07 mesmo sem
    nenhum evento dela ter ocorrido nessa janela.

    Retorna dict no mesmo formato de coletar_ocorrencias_semana, pronto
    para ser mesclado com mesclar_grupos().
    """
    cliente_norm = _norm(cliente)
    grupos = {}
    min_cols = ATIV_COL_STATUSOS + 1

    for row in todos_valores[1:]:
        if len(row) <= min_cols:
            row = row + [""] * (min_cols + 1 - len(row))

        if cliente_norm not in _norm(row[ATIV_COL_CLIENTE]):
            continue

        dt_abertura = _parse_data(row[ATIV_COL_DATA_CRIACAO])
        concluida = _atividade_concluida(row[ATIV_COL_STATUSOS])
        if concluida:
            # Prioriza a data real do Histórico sobre a Data de Conclusão
            # carimbada pela auditoria (que reflete quando NÓS detectamos o
            # "Finalizada", não quando aconteceu de fato — ver
            # _ultima_data_historico). Só cai para a Data de Conclusão se
            # não houver nada de útil no Histórico. Só faz sentido calcular
            # isso quando a OS está de fato concluída — senão qualquer
            # atualização recente de uma OS ainda em andamento poderia ser
            # confundida com um "fechamento" (bug evitado aqui de propósito).
            dt_fecham = _ultima_data_historico(row[ATIV_COL_HISTORICO]) or _parse_data(row[ATIV_COL_DATA_CONCLUSAO])
        else:
            dt_fecham = None

        fechou_na_semana = dt_fecham and data_inicio <= dt_fecham <= data_fim
        abriu_na_semana = dt_abertura and data_inicio <= dt_abertura <= data_fim

        if not (fechou_na_semana or abriu_na_semana):
            continue

        # categoriza pela descrição (onde aparecem palavras como "desligada",
        # "offline", "religamento" etc.) — o código do equipamento sozinho
        # (ex.: "2C-APG100") nunca bate em nenhum padrão, então cai sempre
        # em OUTRAS OCORRÊNCIAS e perde a categorização real (bug corrigido
        # em 13/07/2026). Usa equipamento+descrição juntos pra manter
        # também os casos que só têm palavra-chave no nome do equipamento.
        texto_categoria = f"{row[ATIV_COL_EQUIP]} {row[ATIV_COL_DESCRICAO]}"
        categoria = _rotulo_categoria(texto_categoria)
        status_display = row[ATIV_COL_STATUSOS].strip() or row[ATIV_COL_STATUS].strip() or "Em aberto"
        d = {
            "usina": row[ATIV_COL_USINA].strip() or "Usina não informada",
            "equipamento": row[ATIV_COL_EQUIP].strip() or "Equipamento não informado",
            "numero_os": row[ATIV_COL_NUMEROOS].strip() or None,
            "falha": row[ATIV_COL_DESCRICAO].strip() or "Sem descrição",
            "status": status_display,
            "dt_abertura": dt_abertura,
            "dt_fechamento": dt_fecham,
            "concluida": concluida,
            "origem": "atividade",
        }
        grupos.setdefault(categoria, {"concluidas": [], "abertas": []})

        if fechou_na_semana or (concluida and abriu_na_semana):
            # fechou dentro da semana, OU abriu E já está "Finalizada" mas
            # a Data de Conclusão ficou vazia/fora da janela (bug histórico
            # corrigido em 13/07/2026 — atividades concluídas não gravavam
            # essa data). Sem este fallback, a atividade some do relatório
            # sem deixar rastro.
            grupos[categoria]["concluidas"].append(d)
        else:
            # abriu dentro da semana mas ainda não fechou — mostra como
            # aberta/em andamento, só nesta janela (não perpetuamente).
            grupos[categoria]["abertas"].append(d)

    return {k: v for k, v in grupos.items() if v["concluidas"] or v["abertas"]}


def mesclar_grupos(*grupos_varios):
    """Mescla dois ou mais dicts no formato de coletar_ocorrencias_semana
    (ex.: Painel de Falhas + Painel de Atividades) em um único dict de
    categorias, somando as listas de concluídas/abertas de cada fonte."""
    resultado = {}
    for grupos in grupos_varios:
        for cat, dados in (grupos or {}).items():
            alvo = resultado.setdefault(cat, {"concluidas": [], "abertas": []})
            alvo["concluidas"].extend(dados.get("concluidas", []))
            alvo["abertas"].extend(dados.get("abertas", []))
    return resultado


def listar_usinas_cliente(todos_valores, cliente):
    """Lista todas as usinas do cliente (histórico completo, não só a semana)."""
    cliente_norm = _norm(cliente)
    usinas = []
    vistos = set()
    for row in todos_valores[1:]:
        if len(row) <= COL_USINA:
            continue
        if cliente_norm not in _norm(row[COL_CLIENTE]):
            continue
        usina = row[COL_USINA].strip()
        chave = _norm(usina)
        if usina and chave not in vistos:
            vistos.add(chave)
            usinas.append(usina)
    return sorted(usinas)


def coletar_chamados_abertos(todos_valores, cliente):
    """Ocorrências do cliente (Painel de Falhas) com número de chamado/ticket
    fabricante válido e ainda não concluídas."""
    cliente_norm = _norm(cliente)
    chamados = []
    for row in todos_valores[1:]:
        if len(row) <= COL_DATA_FECHAMENTO:
            row = row + [""] * (COL_DATA_FECHAMENTO + 1 - len(row))
        if cliente_norm not in _norm(row[COL_CLIENTE]):
            continue
        ticket = row[COL_TICKET].strip()
        if not _ticket_valido(ticket) or _is_concluido(row[COL_STATUS]):
            continue
        chamados.append(row)
    return chamados


# ── Geração de conteúdo (offline, sem IA) — parágrafos com negrito por usina ─

def _linha_ocorrencia(d):
    """Linha resumida: Descrição – OS nº (quando houver). Status: X.
    (sem despejar causa/ação/histórico). Corrigido 22/07/2026: usar o
    número da OS (legível pro cliente) em vez do código interno do
    equipamento sempre que disponível, e traduzir o status bruto."""
    falha = d["falha"]
    if len(falha) > 180:
        falha = falha[:177].rstrip() + "..."
    status = _status_exibicao(d["status"])
    numero_os = d.get("numero_os")
    if numero_os:
        return f"{falha} – OS {numero_os}. Status: {status}."
    equip = d["equipamento"]
    return f"{equip} – {falha}. Status: {status}."


def gerar_paragrafos_categoria(categoria, dados):
    """
    Agrupa as ocorrências/atividades de uma categoria por usina (nome da
    usina em negrito, seguido de um bullet resumido por ocorrência).
    Retorna lista de dicts: {"texto": str, "bold": bool}
    """
    todas = list(dados["concluidas"]) + list(dados["abertas"])
    por_usina = {}
    for d in todas:
        por_usina.setdefault(d["usina"], []).append(d)

    paragrafos = []
    for usina in sorted(por_usina.keys()):
        paragrafos.append({"texto": usina, "bold": True})
        for d in por_usina[usina]:
            paragrafos.append({"texto": "• " + _linha_ocorrencia(d), "bold": False})
    return paragrafos


# ── Desligamentos: sempre com data/hora real (aberto via WhatsApp já traz a
# hora do desligamento; ao normalizar, a hora de conclusão) ─────────────────

def _linha_desligamento_runs(d):
    """
    Runs (mistura de negrito/normal na mesma linha) para uma ocorrência de
    desligamento, incluindo data/hora real de abertura e, se concluída, de
    conclusão. Só cai no aviso em negrito se a data realmente estiver
    ausente na origem (caso excepcional — normalmente já vem preenchida
    pela automação do WhatsApp, ou pela Data de Criação/Conclusão da OS).
    """
    equip = d["equipamento"]
    falha = d["falha"]
    status = _status_exibicao(d["status"])
    concluida = d["concluida"]
    dt_abertura = d["dt_abertura"]
    dt_fecham = d["dt_fechamento"]

    runs = [{"texto": f"• {equip} – {falha}. Status: {status}. ", "bold": False}]

    if dt_abertura:
        runs.append({"texto": f"Desligada em {_fmt_data_hora(dt_abertura)}.", "bold": False})
    else:
        runs.append({"texto": "ACRESCENTAR A DATA E HORA DO DESLIGAMENTO", "bold": True})

    if concluida:
        if dt_fecham:
            runs.append({"texto": f" Religada/concluída em {_fmt_data_hora(dt_fecham)}.", "bold": False})
        else:
            runs.append({"texto": " ACRESCENTAR A DATA E HORA DA CONCLUSÃO", "bold": True})

    return runs


def gerar_paragrafos_desligamentos(dados):
    """Igual a gerar_paragrafos_categoria, mas com data/hora real por ocorrência."""
    todas = list(dados["concluidas"]) + list(dados["abertas"])
    por_usina = {}
    for d in todas:
        por_usina.setdefault(d["usina"], []).append(d)

    paragrafos = []
    for usina in sorted(por_usina.keys()):
        paragrafos.append({"runs": [{"texto": usina, "bold": True}]})
        for d in por_usina[usina]:
            paragrafos.append({"runs": _linha_desligamento_runs(d)})
    return paragrafos


# ── Chamados em aberto: uma linha por chamado, usina em negrito inline ─────

def gerar_paragrafos_chamados_fabricante(chamados):
    """Uma linha por chamado da aba ChamadosFabricante (dashboard novo,
    17/07/2026) — mostra Ativo, Identificação do Equipamento, Ticket/RMA,
    Status e a Nota do dashboard (quando o Fred escreveu alguma no popup
    de detalhe). Recebe a lista JÁ filtrada por cliente e sem os
    concluídos — este código só formata, não filtra."""
    if not chamados:
        return [{"texto": "Nenhum chamado em aberto nesta semana.", "bold": False}]

    por_usina = {}
    for c in chamados:
        usina = (c.get("UFV") or "Usina não informada").strip()
        por_usina.setdefault(usina, []).append(c)

    paragrafos = []
    for usina in sorted(por_usina.keys()):
        paragrafos.append({"runs": [{"texto": usina, "bold": True}]})
        for c in por_usina[usina]:
            ativo = (c.get("Ativo") or "").strip()
            equip = (c.get("Identificação do Equipamento") or "").strip()
            ticket = (c.get("Ticket/RMA") or "").strip()
            status = (c.get("Status") or "Em análise").strip()
            nota = (c.get("NotaDashboard") or "").strip()

            descricao_equip = equip or ativo or "Equipamento não informado"
            if equip and ativo and ativo.lower() not in equip.lower():
                descricao_equip = f"{ativo} — {equip}"

            texto_principal = f"- {descricao_equip} – Ticket {ticket or 's/n'} – {status}."
            runs = [{"texto": texto_principal, "bold": False}]
            if nota:
                runs.append({"texto": f" Obs.: {nota}", "bold": False})
            paragrafos.append({"runs": runs})
    return paragrafos


def gerar_paragrafos_chamados(chamados):
    """Uma linha por chamado válido: '- Usina – Case #ticket -> resumo -> status.'
    (chamados vêm sempre do Painel de Falhas — Painel de Atividades não tem
    número de chamado/ticket de fabricante, só número de OS interno)."""
    if not chamados:
        return [{"texto": "Nenhum chamado em aberto nesta semana.", "bold": False}]

    por_usina = {}
    for row in chamados:
        usina = row[COL_USINA].strip() or "Usina não informada"
        por_usina.setdefault(usina, []).append(row)

    paragrafos = []
    for usina in sorted(por_usina.keys()):
        for row in por_usina[usina]:
            ticket = row[COL_TICKET].strip()
            resumo = row[COL_FALHA].strip() or row[COL_CAUSA].strip() or "Sem descrição detalhada"
            status = row[COL_STATUS].strip() or "Em análise"
            runs = [
                {"texto": "- ", "bold": False},
                {"texto": usina, "bold": True},
                {"texto": f" – Case #{ticket} -> {resumo} -> {status}.", "bold": False},
            ]
            paragrafos.append({"runs": runs})
    return paragrafos


# Layout real da aba "Zeladoria" (gid 987654321): linha 1 = título dos grupos,
# linha 2 = subcabeçalhos, dados a partir da linha 3.
# Colunas (0-based): A Cliente, B Usina, depois 4 grupos de 4 colunas cada
# (Última Data, Próxima Data, Quantidade, Status): Roçada, Poda Química, Lavagem dos Módulos, Controle de Pragas.
ZEL_COL_CLIENTE, ZEL_COL_USINA = 0, 1
ZEL_GRUPOS = [
    ("Roçada",               2),
    ("Poda Química",         6),
    ("Lavagem dos Módulos", 10),
    ("Controle de Pragas",  14),
]


# Clientes onde uma célula de Zeladoria vazia deve exibir "Acompanhamento"
# em vez de "Sem informação" — combinado com Fred em 22/07/2026: ele ajusta
# manualmente quando necessário, mas o padrão nesses dois clientes deve
# partir de "Acompanhamento".
CLIENTES_ZELADORIA_DEFAULT_ACOMPANHAMENTO = ["gd energy", "alves lima"]


def coletar_zeladoria(zeladoria_valores, cliente):
    """
    zeladoria_valores: ws.get_all_values() da aba Zeladoria (linhas 1-2 = cabeçalho).
    Retorna lista de dicts: [{"usina": ..., "grupos": [{"nome":, "status":, "ultima_data":}, ...]}, ...]
    """
    cliente_norm = _norm(cliente)
    usa_default_acompanhamento = any(c in cliente_norm for c in CLIENTES_ZELADORIA_DEFAULT_ACOMPANHAMENTO)
    resultado = []
    for row in zeladoria_valores[2:]:
        if len(row) <= ZEL_GRUPOS[-1][1] + 3:
            row = row + [""] * (ZEL_GRUPOS[-1][1] + 4 - len(row))
        if not row[ZEL_COL_USINA].strip():
            continue
        if cliente_norm not in _norm(row[ZEL_COL_CLIENTE]):
            continue

        grupos_usina = []
        for nome, col_ini in ZEL_GRUPOS:
            ultima_data = row[col_ini].strip()
            status = row[col_ini + 3].strip()
            if not status and usa_default_acompanhamento:
                status = "Acompanhamento"
            grupos_usina.append({"nome": nome, "status": status, "ultima_data": ultima_data})

        resultado.append({"usina": row[ZEL_COL_USINA].strip(), "grupos": grupos_usina})
    return resultado


def _status_zeladoria_para_tabela(item, nome_grupo):
    """Texto de uma célula da tabela de Zeladoria (só o status, conforme
    pedido — sem despejar a última data dentro da célula para não poluir)."""
    grupo = next((g for g in item["grupos"] if g["nome"] == nome_grupo), None)
    if not grupo:
        return "Sem informação"
    return grupo["status"] or "Sem informação"


def paginar_zeladoria(zeladoria_usinas, max_por_pagina=ZELADORIA_MAX_USINAS_POR_PAGINA):
    """Divide as usinas da Zeladoria em páginas de tabela, pra nunca estourar
    o slide (era o problema antigo: texto cortado com muitas usinas)."""
    if not zeladoria_usinas:
        return [[]]
    return [zeladoria_usinas[i:i + max_por_pagina] for i in range(0, len(zeladoria_usinas), max_por_pagina)]


# ── Clonagem de slides (mantém 100% do estilo do modelo Grid Co.) ──────────

def _duplicate_slide(prs, index):
    """Duplica um slide existente do modelo (mesmo layout, cores, fontes, formas)."""
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)

    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    id_map = {}
    for rId, rel in source.part.rels.items():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        new_rid = dest.part.rels._add_relationship(rel.reltype, rel._target, rel.is_external)
        id_map[rId] = new_rid

    for shp in source.shapes:
        newel = copy.deepcopy(shp._element)
        for el in newel.iter():
            for attr in ("embed", "link", "id"):
                full = qn(f"r:{attr}")
                val = el.get(full)
                if val and val in id_map:
                    el.set(full, id_map[val])
        dest.shapes._spTree.append(newel)
    return dest


def _extrair_formato_base(shape):
    """Extrai fonte/tamanho/itálico/cor do 1º run existente na shape (pra reaplicar depois)."""
    tf = shape.text_frame
    primeiro_par = tf.paragraphs[0]
    if not primeiro_par.runs:
        return None
    font = primeiro_par.runs[0].font
    cor_rgb = None
    try:
        if font.color and font.color.type is not None:
            cor_rgb = font.color.rgb  # levanta AttributeError se for cor de tema (schemeClr)
    except AttributeError:
        cor_rgb = None
    return (font.name, font.size, font.italic, cor_rgb)


def _extrair_ppr_bullet(shape):
    """
    Extrai o XML de formatação de parágrafo (<a:pPr>, que carrega a definição
    do bullet/marcador ">" do modelo) do 1º parágrafo da shape, ANTES de
    qualquer tf.clear(). Sem isso, só a 1ª linha herdava o bullet do modelo —
    as linhas adicionadas via add_paragraph() vinham sem marcador nenhum.
    """
    tf = shape.text_frame
    p_origem = tf.paragraphs[0]._p
    pPr = p_origem.find(qn("a:pPr"))
    return copy.deepcopy(pPr) if pPr is not None else None


def _aplicar_ppr_bullet(paragrafo, ppr_modelo):
    """Aplica (substituindo) o <a:pPr> de ppr_modelo num parágrafo já criado."""
    if ppr_modelo is None:
        return
    p_destino = paragrafo._p
    pPr_atual = p_destino.find(qn("a:pPr"))
    novo = copy.deepcopy(ppr_modelo)
    if pPr_atual is not None:
        p_destino.remove(pPr_atual)
    p_destino.insert(0, novo)


def _set_text_preservando_estilo(shape, novo_texto):
    """Substitui o texto de um text box mantendo a formatação (incl. negrito) do 1º run
    E o marcador de bullet (">") do modelo em TODAS as linhas, não só na 1ª."""
    tf = shape.text_frame
    primeiro_par = tf.paragraphs[0]
    bold_original = primeiro_par.runs[0].font.bold if primeiro_par.runs else None
    fmt = _extrair_formato_base(shape)
    ppr_bullet = _extrair_ppr_bullet(shape)

    tf.clear()
    linhas = novo_texto.split("\n")
    for i, linha in enumerate(linhas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _aplicar_ppr_bullet(p, ppr_bullet)
        run = p.add_run()
        run.text = linha
        if fmt:
            if fmt[0]:
                run.font.name = fmt[0]
            if fmt[1]:
                run.font.size = fmt[1]
            if fmt[2] is not None:
                run.font.italic = fmt[2]
            if fmt[3]:
                run.font.color.rgb = fmt[3]
        if bold_original is not None:
            run.font.bold = bold_original


def _set_paragrafos_com_estilo(shape, paragrafos):
    """
    Substitui o conteúdo de um text box por uma lista de parágrafos, permitindo
    alternar negrito por parágrafo (usado para destacar nome de usina/categoria)
    OU, quando o item tiver a chave "runs", misturar negrito/normal dentro da
    MESMA linha (usado em Desligamentos e Chamados em Aberto).

    paragrafos: lista de dicts, em um dos dois formatos:
      - {"texto": str, "bold": bool}                      -> parágrafo simples
      - {"runs": [{"texto": str, "bold": bool}, ...]}      -> parágrafo com runs mistos
    """
    fmt = _extrair_formato_base(shape)
    ppr_bullet = _extrair_ppr_bullet(shape)
    tf = shape.text_frame
    tf.clear()

    if not paragrafos:
        paragrafos = [{"texto": "", "bold": False}]

    for i, p in enumerate(paragrafos):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _aplicar_ppr_bullet(par, ppr_bullet)
        runs_def = p.get("runs") or [{"texto": p.get("texto", ""), "bold": p.get("bold", False)}]
        for r in runs_def:
            run = par.add_run()
            run.text = r.get("texto", "")
            if fmt:
                if fmt[0]:
                    run.font.name = fmt[0]
                if fmt[1]:
                    run.font.size = fmt[1]
                if fmt[2] is not None:
                    run.font.italic = fmt[2]
                if fmt[3]:
                    run.font.color.rgb = fmt[3]
            run.font.bold = bool(r.get("bold", False))


def _formatar_celula_tabela(celula, fmt, negrito):
    """Aplica a fonte/tamanho extraídos do template (fmt) numa célula de tabela."""
    celula.text_frame.word_wrap = True
    for p in celula.text_frame.paragraphs:
        p.alignment = p.alignment  # no-op, mantém padrão do template de tabela
        if not p.runs:
            p.add_run()
        for r in p.runs:
            if fmt:
                if fmt[0]:
                    r.font.name = fmt[0]
                if fmt[1]:
                    r.font.size = fmt[1]
            r.font.bold = negrito


def _remover_fotos(slide):
    """Remove os placeholders de foto (o gerador não tem fotos de campo)."""
    for shp in list(slide.shapes):
        if shp.shape_type == 13 and (shp.name or "").lower().startswith("imagem"):
            shp._element.getparent().remove(shp._element)


def _find_shape(slide, contem_texto):
    for shp in slide.shapes:
        if shp.has_text_frame and contem_texto.lower() in shp.text_frame.text.lower():
            return shp
    return None


# ── Montagem do PPTX final ──────────────────────────────────────────────────

def _extrair_orcamento_linhas(shape, fmt):
    """
    Estima quantas linhas cabem na caixa do modelo e quantos caracteres cabem
    por linha, a partir da altura/largura reais da shape (EMU) e do tamanho
    de fonte do template — substitui o orçamento por "número de caracteres",
    que não refletia a quebra de linha real e deixava o conteúdo estourar
    a página em categorias com texto mais longo.
    """
    altura_pt = shape.height / 12700
    largura_pt = shape.width / 12700
    tamanho_fonte_pt = fmt[1].pt if fmt and fmt[1] else 18
    altura_linha_pt = tamanho_fonte_pt * 1.35  # inclui folga de espaçamento entre parágrafos
    largura_media_char_pt = tamanho_fonte_pt * 0.52  # estimativa conservadora (fonte não é monoespaçada)
    max_linhas = max(1, int(altura_pt // altura_linha_pt))
    chars_por_linha = max(15, int(largura_pt // largura_media_char_pt))
    max_linhas = max(1, int(max_linhas * 0.8))  # margem de segurança de 20%
    return max_linhas, chars_por_linha


def _texto_paragrafo(paragrafo):
    if "runs" in paragrafo:
        return "".join(r.get("texto", "") for r in paragrafo["runs"])
    return paragrafo.get("texto", "")


def _linhas_ocupadas(paragrafo, chars_por_linha):
    texto = _texto_paragrafo(paragrafo)
    if not texto:
        return 1
    return max(1, -(-len(texto) // chars_por_linha))  # divisão com arredondamento pra cima


def _linhas_totais(paragrafos, chars_por_linha):
    return sum(_linhas_ocupadas(p, chars_por_linha) for p in paragrafos)


def gerar_paragrafos_multi_categoria(dados_por_categoria):
    """Igual a gerar_paragrafos_categoria, mas mesclando várias categorias
    (ex.: Comunicações + Controle de Vegetação + Outras Ocorrências) numa
    única lista agrupada só por usina, sem subtítulo de categoria — usado
    quando a pauta do relatório não distingue esses sub-tópicos e tudo cai
    sob um único tópico "ATIVIDADES SEMANAIS" (combinado 22/07/2026: a
    pauta fixa não tem "Controle de Vegetação" nem "Outras Ocorrências"
    como itens separados, então o conteúdo não pode ficar em páginas com
    esses títulos)."""
    por_usina = {}
    for dados in dados_por_categoria.values():
        todas = list(dados["concluidas"]) + list(dados["abertas"])
        for d in todas:
            por_usina.setdefault(d["usina"], []).append(d)

    paragrafos = []
    for usina in sorted(por_usina.keys()):
        paragrafos.append({"texto": usina, "bold": True})
        for d in por_usina[usina]:
            paragrafos.append({"texto": "• " + _linha_ocorrencia(d), "bold": False})
    return paragrafos


def paginar_paragrafos_simples(paragrafos, max_linhas, chars_por_linha):
    """Pagina uma lista plana de parágrafos (sem cabeçalho de categoria) por
    orçamento de linhas — usado pela seção mesclada ATIVIDADES SEMANAIS."""
    paginas, atual, atual_linhas = [], [], 0
    for p in paragrafos:
        linhas = _linhas_ocupadas(p, chars_por_linha)
        if atual and (atual_linhas + linhas > max_linhas):
            paginas.append(atual)
            atual, atual_linhas = [], 0
        atual.append(p)
        atual_linhas += linhas
    if atual:
        paginas.append(atual)
    return paginas or [[]]


def agrupar_em_paginas(grupos, max_linhas, chars_por_linha, ordem=None):
    """
    Agrupa categorias em páginas por orçamento de LINHAS (estimadas a partir
    da caixa real do modelo). Retorna lista de páginas, cada página = lista
    de tuplas (categoria, paragrafos).
    `ordem`, se fornecida, define a sequência das categorias (categorias fora
    dela vão para o final, em ordem alfabética).
    """
    chaves = _ordenar_categorias(grupos.keys()) if ordem is None else ordem
    blocos = [(cat, gerar_paragrafos_categoria(cat, grupos[cat])) for cat in chaves if cat in grupos]
    paginas, atual, atual_linhas = [], [], 0
    for cat, paragrafos in blocos:
        linhas_bloco = 1 + _linhas_totais(paragrafos, chars_por_linha)  # +1 = cabeçalho da categoria
        if atual and (atual_linhas + linhas_bloco > max_linhas):
            paginas.append(atual)
            atual, atual_linhas = [], 0
        atual.append((cat, paragrafos))
        atual_linhas += linhas_bloco
    if atual:
        paginas.append(atual)
    return paginas


def agrupar_desligamentos_em_paginas(dados, max_linhas, chars_por_linha):
    """Pagina os parágrafos (com runs) de Desligamentos por orçamento de linhas."""
    paragrafos = gerar_paragrafos_desligamentos(dados)
    paginas, atual, atual_linhas = [], [], 0
    for p in paragrafos:
        linhas = _linhas_ocupadas(p, chars_por_linha)
        if atual and (atual_linhas + linhas > max_linhas):
            paginas.append(atual)
            atual, atual_linhas = [], 0
        atual.append(p)
        atual_linhas += linhas
    if atual:
        paginas.append(atual)
    return paginas or [[]]


def _renderizar_pagina_categoria(prs, titulo, corpo_paragrafos):
    """Cria um slide de conteúdo (clonado do modelo) com título + corpo em texto."""
    novo = _duplicate_slide(prs, 2)
    _remover_fotos(novo)
    shp_categoria = _find_shape(novo, "DESLIGAMENTOS")
    shp_corpo = _find_shape(novo, "Foram registradas")
    if shp_categoria:
        _set_text_preservando_estilo(shp_categoria, titulo)
    if shp_corpo:
        _set_paragrafos_com_estilo(shp_corpo, corpo_paragrafos)
    return novo


def _renderizar_pagina_zeladoria_tabela(prs, titulo, pagina_usinas):
    """
    Cria um slide de Zeladoria com uma TABELA real (Usina | Roçagem | Poda
    Química | Limpeza dos Módulos), no lugar da antiga lista de bullets —
    o texto corrido não cabia numa única página quando havia muitas usinas.
    """
    novo = _duplicate_slide(prs, 2)
    _remover_fotos(novo)
    shp_categoria = _find_shape(novo, "DESLIGAMENTOS")
    shp_corpo = _find_shape(novo, "Foram registradas")
    if shp_categoria:
        _set_text_preservando_estilo(shp_categoria, titulo)

    if shp_corpo:
        fmt = _extrair_formato_base(shp_corpo)
        left, top, width, height = shp_corpo.left, shp_corpo.top, shp_corpo.width, shp_corpo.height
        shp_corpo._element.getparent().remove(shp_corpo._element)
    else:
        fmt = None
        left = top = width = height = None

    if left is None:
        # fallback conservador, caso o placeholder de corpo não seja encontrado
        from pptx.util import Emu
        left, top, width, height = Emu(700000), Emu(1500000), Emu(11500000), Emu(5500000)

    linhas = len(pagina_usinas) + 1  # +1 cabeçalho
    colunas = len(ZELADORIA_COLUNAS)
    graphic_frame = novo.shapes.add_table(linhas, colunas, left, top, width, height)
    tabela = graphic_frame.table

    for j, titulo_col in enumerate(ZELADORIA_COLUNAS):
        cel = tabela.cell(0, j)
        cel.text = titulo_col
        _formatar_celula_tabela(cel, fmt, negrito=True)

    for i, item in enumerate(pagina_usinas, start=1):
        cel_usina = tabela.cell(i, 0)
        cel_usina.text = item["usina"]
        _formatar_celula_tabela(cel_usina, fmt, negrito=True)
        for j, (_, nome_grupo) in enumerate(ZELADORIA_MAPA_COLUNAS, start=1):
            cel = tabela.cell(i, j)
            cel.text = _status_zeladoria_para_tabela(item, nome_grupo)
            _formatar_celula_tabela(cel, fmt, negrito=False)

    return novo


def gerar_relatorio_pptx(cliente, semana_num, data_label, grupos, chamados=None, zeladoria_usinas=None, usinas_cliente=None):
    """
    grupos: dict categoria -> {"concluidas": [dict,...], "abertas": [dict,...]}
            (normalmente o resultado de mesclar_grupos(coletar_ocorrencias_semana(...),
             coletar_atividades_semana(...)) — ocorrências do Painel de Falhas E
             OSs do Painel de Atividades, já combinadas por categoria).
    semana_num: número da semana do ano (baseado na data final do período)
    data_label: data final formatada (dd/mm/aaaa) — mantido por compatibilidade
                de assinatura, não é mais exibido na capa (padrão Grid Co.).
    chamados: lista de linhas (Painel de Falhas) com chamado/ticket válido em aberto (opcional)
    zeladoria_usinas: retorno de coletar_zeladoria() (opcional)
    usinas_cliente: lista de nomes de usinas do cliente (ex.: listar_usinas_cliente()) —
                     usada para montar a página "OUTROS TEMAS" (uma seção em branco por
                     usina, pronta para o Fred preencher manualmente). Opcional; se omitida
                     ou vazia, a página "OUTROS TEMAS" não é criada.
    Retorna BytesIO() pronto para download.
    """
    prs = Presentation(TEMPLATE_PATH)

    # Orçamento de paginação calculado UMA VEZ a partir das dimensões reais
    # do placeholder de corpo do modelo (slide 2), pra nunca estourar a página.
    _corpo_modelo = _find_shape(prs.slides[2], "Foram registradas")
    if _corpo_modelo is not None:
        _fmt_corpo_modelo = _extrair_formato_base(_corpo_modelo)
        max_linhas, chars_por_linha = _extrair_orcamento_linhas(_corpo_modelo, _fmt_corpo_modelo)
    else:
        max_linhas, chars_por_linha = 18, 70  # fallback conservador

    grupos = dict(grupos)  # não alterar o dict do chamador
    dados_comunicacoes = grupos.pop(CAT_COMUNICACOES, None)
    dados_desligamentos = grupos.pop(CAT_DESLIGAMENTOS, None)
    outras_categorias = grupos  # Inversores, Strings/Módulos, etc. — ocorrências E OSs

    # --- Slide 1: Capa — só cliente e semana (sem data) ---------------------
    capa = _duplicate_slide(prs, 0)
    shp_titulo = _find_shape(capa, "O&M")
    if shp_titulo:
        _set_text_preservando_estilo(shp_titulo, f"O&M – {cliente}\nSemana {semana_num}")

    # --- Slide 2: Ata da reunião ---------------------------------------------
    # Pauta padrão fixa para todos os clientes, exceto RenoGrid (que mantém a
    # pauta dinâmica, montada a partir do que realmente aparece na semana).
    if _norm(cliente) in CLIENTES_PAUTA_DINAMICA:
        topicos = []
        if dados_comunicacoes:
            topicos.append(CAT_COMUNICACOES)
        if outras_categorias:
            topicos.append("OCORRÊNCIAS SEMANAIS")
        if dados_desligamentos:
            topicos.append(CAT_DESLIGAMENTOS)
        if chamados:
            topicos.append("CHAMADOS EM ABERTO")
        if zeladoria_usinas:
            topicos.append("ZELADORIA")
    else:
        topicos = list(PAUTA_PADRAO_ATA)
    ata = _duplicate_slide(prs, 1)
    shp_lista = _find_shape(ata, "Desligamentos")
    if shp_lista:
        _set_text_preservando_estilo(shp_lista, "\n".join(topicos))

    pauta_fixa = _norm(cliente) not in CLIENTES_PAUTA_DINAMICA

    # --- Página(s): Comunicações / SCADA / CCTV — sempre primeiro, própria(s)
    #     página(s), separada de Inversores/Strings/demais equipamentos -------
    #     (só nos clientes com pauta dinâmica — RenoGrid; nos demais, entra
    #     mesclada em ATIVIDADES SEMANAIS mais abaixo)
    if dados_comunicacoes and not pauta_fixa:
        for pagina in agrupar_em_paginas({CAT_COMUNICACOES: dados_comunicacoes}, max_linhas, chars_por_linha):
            cat, paragrafos = pagina[0]
            _renderizar_pagina_categoria(prs, cat, paragrafos)

    if pauta_fixa:
        # --- Página(s): ATIVIDADES SEMANAIS — mescla Comunicações + todas as
        #     demais categorias (vegetação, ocorrências etc.) num único
        #     tópico, já que a pauta fixa não distingue esses sub-tópicos.
        dados_mesclados = dict(outras_categorias)
        if dados_comunicacoes:
            dados_mesclados[CAT_COMUNICACOES] = dados_comunicacoes
        if dados_mesclados:
            paragrafos_mesclados = gerar_paragrafos_multi_categoria(dados_mesclados)
            paginas_ativ = paginar_paragrafos_simples(paragrafos_mesclados, max_linhas, chars_por_linha)
            for i, pagina in enumerate(paginas_ativ):
                titulo = "ATIVIDADES SEMANAIS" if i == 0 else "ATIVIDADES SEMANAIS (cont.)"
                _renderizar_pagina_categoria(prs, titulo, pagina)
    else:
        # --- Página(s): "Ocorrências da Semana" — Inversores, Strings/Módulos e
        #     demais equipamentos, já com Painel de Falhas + Painel de Atividades
        #     combinados por categoria (empacotadas por orçamento de linhas) -----
        for pagina in agrupar_em_paginas(outras_categorias, max_linhas, chars_por_linha):
            if len(pagina) == 1:
                categoria, paragrafos = pagina[0]
                _renderizar_pagina_categoria(prs, categoria, paragrafos)
            else:
                combinados = []
                for cat, paragrafos in pagina:
                    combinados.append({"texto": cat, "bold": True})
                    combinados.extend(paragrafos)
                _renderizar_pagina_categoria(prs, "OCORRÊNCIAS DA SEMANA", combinados)

    # --- Página(s): Desligamentos — com data/hora real de abertura/conclusão.
    #     Nos clientes de pauta fixa, "DESLIGAMENTOS" é sempre um item da
    #     pauta — a página é sempre gerada, mesmo sem nenhum desligamento no
    #     período (mesmo padrão da página OUTROS TEMAS: em branco, por
    #     usina, pronta pra preenchimento manual se necessário).
    if dados_desligamentos:
        paginas_deslig = agrupar_desligamentos_em_paginas(dados_desligamentos, max_linhas, chars_por_linha)
        for i, pagina in enumerate(paginas_deslig):
            titulo = CAT_DESLIGAMENTOS if i == 0 else f"{CAT_DESLIGAMENTOS} (cont.)"
            _renderizar_pagina_categoria(prs, titulo, pagina)
    elif pauta_fixa and usinas_cliente:
        corpo_deslig_vazio = [{"texto": usina, "bold": True} for usina in usinas_cliente]
        _renderizar_pagina_categoria(prs, CAT_DESLIGAMENTOS, corpo_deslig_vazio)

    # --- Página: Chamados em aberto (aba ChamadosFabricante, dashboard) ----
    if chamados:
        _renderizar_pagina_categoria(prs, "CHAMADOS EM ABERTO", gerar_paragrafos_chamados_fabricante(chamados))

    # --- Página(s): Zeladoria — tabela real (Usina | Roçagem | Poda Química |
    #     Limpeza dos Módulos), paginada para nunca estourar o slide ---------
    if zeladoria_usinas:
        for i, pagina in enumerate(paginar_zeladoria(zeladoria_usinas)):
            titulo = "ZELADORIA" if i == 0 else "ZELADORIA (cont.)"
            _renderizar_pagina_zeladoria_tabela(prs, titulo, pagina)

    # --- Página: Outros Temas — em branco, uma seção em negrito por usina do
    #     cliente, para o Fred preencher manualmente depois (ex.: pendências
    #     com seguradora, PV Operation etc.). Combinado 22/07/2026: sempre
    #     gerada (não depende de nenhuma fonte de dados automática).
    if usinas_cliente:
        corpo_outros_temas = []
        for usina in usinas_cliente:
            corpo_outros_temas.append({"texto": usina, "bold": True})
        _renderizar_pagina_categoria(prs, "OUTROS TEMAS", corpo_outros_temas)

    # --- Reordena o deck: capa nova -> ata nova -> conteúdo -> contato -----
    xml_slides = prs.slides._sldIdLst
    todos_els = list(xml_slides)
    contato_el = todos_els[6]
    capa_el = todos_els[7]
    ata_el = todos_els[8]
    conteudo_els = todos_els[9:]

    for e in todos_els:
        xml_slides.remove(e)
    for e in [capa_el, ata_el] + conteudo_els + [contato_el]:
        xml_slides.append(e)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
