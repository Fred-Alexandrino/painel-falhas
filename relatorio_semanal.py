# -*- coding: utf-8 -*-
"""
relatorio_semanal.py
─────────────────────────────────────────────────────────────────────────
Geração automática do "Relatório Semanal" (.pptx) no modelo Grid Co.,
a partir das ocorrências (Painel de Falhas) E das OSs/atividades (Painel
de Atividades) já registradas nas planilhas.

Não usa nenhuma API de IA — o texto é montado por regras determinísticas
em cima dos campos de cada fonte (resumido, sem despejar o histórico
bruto de ações).
"""

import os
import re
import copy
import unicodedata
from io import BytesIO
from datetime import datetime

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.util import Pt
from pptx.dml.color import RGBColor

# ── Configuração ───────────────────────────────────────────────────────────

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "templates", "modelo_relatorio_semanal.pptx")

# Índices das colunas da planilha principal — Painel de Falhas (0-based,
# iguais ao app.py / CAMPO_COL)
COL_ID, COL_CLIENTE, COL_USINA, COL_EQUIP = 0, 1, 2, 3
COL_FALHA, COL_CAUSA, COL_IMPACTADOS, COL_ACAO = 4, 5, 6, 7
COL_STATUS, COL_TICKET, COL_OS, COL_HISTORICO = 8, 9, 10, 11
COL_DATA_ABERTURA = 12
COL_DATA_FECHAMENTO = 20  # coluna U (21ª, 1-based) = índice 20 (0-based)

# Índices das colunas da aba Painel de Atividades (0-based, iguais ao
# ATIV_CAMPO_COL do app.py, que é 1-based -> aqui -1)
ATIV_COL_ID, ATIV_COL_CLIENTE, ATIV_COL_USINA, ATIV_COL_EQUIP = 0, 1, 2, 3
ATIV_COL_DESCRICAO, ATIV_COL_RESPONSAVEL, ATIV_COL_PRAZO, ATIV_COL_PRIORIDADE = 4, 5, 6, 7
ATIV_COL_STATUS, ATIV_COL_DATA_CRIACAO, ATIV_COL_DATA_CONCLUSAO, ATIV_COL_HISTORICO = 8, 9, 10, 11
ATIV_COL_EDITOR, ATIV_COL_NUMEROOS, ATIV_COL_STATUSOS = 12, 13, 14

STATUS_CONCLUIDO = ["concluído", "concluido", "resolvido", "fechado", "resolved", "closed"]

# CORRIGIDO 22/07/2026: uma atividade/OS só conta como "realizada" no
# RELATÓRIO SEMANAL (cliente-facing) quando o statusOS chega em
# "Finalizada" — igual à aba "Concluídas" do próprio Painel de Atividades
# (statusOS=="Finalizada"; "Em Revisão" é a aba separada "Em Verificação").
# A versão anterior também tratava "Em Revisão" como concluída aqui, o que
# fazia o relatório mostrar como pronta uma atividade ainda em validação
# (a fallback de "abriu E já está 'concluída'" pegava OSs criadas na
# semana que só tinham chegado em Em Revisão) — Fred precisou corrigir
# manualmente o relatório da Alves Lima (Semana 30) por causa disso.
# Esta definição é local a este módulo — não afeta as abas/KPIs do
# dashboard em app.py (que usam sua própria lógica de Em Verificação).
STATUS_OS_CONCLUIDO = ["finalizada"]

# Falhas já convertidas em Atividade não devem aparecer de novo como
# ocorrência "em aberto" no relatório — o trabalho já é reportado (com o
# número de OS) através de coletar_atividades_semana(). Sem esta exclusão,
# como "Convertida em Atividade" nunca bate em STATUS_CONCLUIDO, a
# ocorrência cai perpetuamente em "aberta"/backlog e aparece em toda
# semana seguinte, duplicando a mesma informação. Corrigido 22/07/2026.
STATUS_FALHA_EXCLUIR_RELATORIO = ["convertida em atividade", "convertido em atividade"]

CAT_COMUNICACOES = "COMUNICAÇÕES / SCADA / CCTV"
CAT_DESLIGAMENTOS = "DESLIGAMENTOS"

# Pauta padrão da "Ata da reunião" (slide 2) — combinado com Fred em
# 22/07/2026: para todo cliente, EXCETO RenoGrid, a pauta é sempre esta
# lista fixa, independente das categorias que de fato tiveram ocorrência
# na semana. RenoGrid mantém a pauta dinâmica (montada a partir do que
# realmente aparece no relatório daquela semana).
PAUTA_PADRAO_ATA = ["ATIVIDADES SEMANAIS", CAT_DESLIGAMENTOS, "PCM"]
CLIENTES_PAUTA_DINAMICA = ["renogrid"]

# Categorias "bonitas" para exibir no relatório (padrões regex, avaliados em ordem —
# o primeiro que bater define a categoria). Casa tanto nomes por extenso ("Inversor")
# quanto códigos curtos usados em campo ("Inv-18", "NVR", "Chave Seccionadora Skid 1").
# A ordem desta lista também define a ordem de exibição das categorias no relatório
# (exceto Comunicações e Desligamentos, que sempre recebem página própria — ver
# gerar_relatorio_pptx).
PADROES_CATEGORIA = [
    (r"\binv[\s\-]*r?\.?\d+|inversor", "INVERSORES"),
    (r"tracker|\btcu\b", "TRACKERS / TCU"),
    (r"fusive|fusíve", "FUSÍVEIS"),
    (r"transformador", "TRANSFORMADORES"),
    (r"switchgear|chave seccionadora|disjuntor|religador", "SWITCHGEAR"),
    (r"\bskid", "SKIDS"),
    (r"preventiv", "MANUTENÇÃO PREVENTIVA"),
    (r"\bstring\b|modulo|módulo|otimizador", "STRINGS / MÓDULOS"),
    (r"usina desligad|usina offline|desligamento|ufv desligad|usina parad", CAT_DESLIGAMENTOS),
    (r"comunica|scada|cctv|cftv|\bnvr\b|camera|câmera|speed dome|igate", CAT_COMUNICACOES),
    (r"sensor|piranometro|piranômetro|solarimetric|estacao solarimetrica|estação solarimétrica", "SENSORES"),
    (r"emergenc", "EMERGÊNCIAS"),
    (r"operac|opera[çc][aã]o", "OPERAÇÕES"),
    (r"termografia|termograf", "TERMOGRAFIA"),
    (r"restart|religamento", "RESTART / RELIGAMENTO"),
    (r"ronda", "RONDAS SEMANAIS"),
    (r"exaustor", "EXAUSTORES"),
    (r"vegeta", "CONTROLE DE VEGETAÇÃO"),
]

# Rótulo genérico pra tudo que não bate em nenhum padrão acima (formulários,
# spare parts, reclamações de concessionária, códigos Fracttal sem nenhuma
# palavra reconhecível etc.) — nunca usa o texto bruto do equipamento como
# categoria própria, pra não fragmentar o relatório em uma mini-seção por
# código.
CAT_OUTRAS = "OUTRAS OCORRÊNCIAS"

# Ordem de exibição das categorias "genéricas" (tudo que não é Comunicações
# nem Desligamentos, que são tratadas à parte). Categorias fora desta lista
# (rótulos livres vindos direto do campo Equipamento) vão para o final, em
# ordem alfabética.
ORDEM_CATEGORIAS_GERAIS = [rotulo for _, rotulo in PADROES_CATEGORIA
                           if rotulo not in (CAT_COMUNICACOES, CAT_DESLIGAMENTOS)] + [CAT_OUTRAS]

# Zeladoria: colunas fixas da tabela do relatório
ZELADORIA_COLUNAS = ["Usina", "Roçagem", "Poda Química", "Limpeza dos Módulos"]
# Mapeia cada coluna da tabela para o "nome" de grupo usado em coletar_zeladoria()
ZELADORIA_MAPA_COLUNAS = [("Roçagem", "Roçada"), ("Poda Química", "Poda Química"),
                          ("Limpeza dos Módulos", "Lavagem dos Módulos")]
ZELADORIA_MAX_USINAS_POR_PAGINA = 8  # evita estourar a tabela pra fora do slide


# ── Utilidades de texto/data ────────────────────────────────────────────────

def _norm(s):
    s = unicodedata.normalize("NFKD", (s or "").lower())
    s = s.encode("ascii", "ignore").decode("ascii")
    return s.strip()


def _rotulo_categoria(equip_bruto):
    n = _norm(equip_bruto)
    for padrao, rotulo in PADROES_CATEGORIA:
        if re.search(padrao, n):
            return rotulo
    return CAT_OUTRAS


def _parse_data(txt):
    """Aceita 'dd/mm/aaaa[ hh:mm[:ss]]' -> datetime, ou None."""
    if not txt:
        return None
    txt = txt.strip()
    for fmt in ("%d/%m/%Y %H:%M:%S", "%d/%m/%Y %H:%M", "%d/%m/%Y"):
        try:
            return datetime.strptime(txt, fmt)
        except ValueError:
            continue
    return None


_RE_DATA_HISTORICO = re.compile(r"(\d{2}/\d{2}/\d{4})\s+(\d{2}:\d{2}(?::\d{2})?)")
_RE_MARCA_CONCLUSAO_HISTORICO = re.compile(r"100%|para [\"']em revis|para [\"']finalizad", re.IGNORECASE)


def _ultima_data_historico(historico_texto):
    """
    Extrai do Histórico da OS a data/hora em que o trabalho de fato foi
    concluído — a PRIMEIRA linha que menciona "100%" ou uma transição de
    status para "Em Revisão"/"Finalizada" (o marco real de "terminei o
    trabalho em campo"). NÃO usa a última linha do histórico: linhas
    posteriores costumam ser só confirmações administrativas tardias (ex.:
    a Fracttal só oficializar "Finalizada" dias depois), que têm a mesma
    armadilha da Data de Conclusão original — refletem quando ALGUÉM
    confirmou, não quando o trabalho realmente terminou.

    Caso relatado pelo Fred: OS 9173 chegou a 100% em 10/07 (linha do
    histórico: "...progresso da tarefa foi de 0% para 100%."), mas seguiu
    aparecendo em relatórios de semanas seguintes porque uma versão
    anterior desta função pegava a última data do histórico (que podia ser
    uma confirmação tardia), e antes disso porque a Data de Conclusão era
    carimbada só quando a auditoria detectava a mudança. Corrigido em
    22/07/2026 (ver histórico de correções — este é o segundo ajuste para
    o mesmo caso).

    Se nenhuma linha com esses marcadores existir, cai para a última data
    encontrada no histórico (comportamento anterior, como fallback).
    """
    if not historico_texto:
        return None
    primeira_marca = None
    ultima_qualquer = None
    for linha in historico_texto.splitlines():
        m = _RE_DATA_HISTORICO.search(linha)
        if not m:
            continue
        dt = _parse_data(f"{m.group(1)} {m.group(2)}")
        if not dt:
            continue
        if ultima_qualquer is None or dt > ultima_qualquer:
            ultima_qualquer = dt
        if _RE_MARCA_CONCLUSAO_HISTORICO.search(linha) and (primeira_marca is None or dt < primeira_marca):
            primeira_marca = dt
    return primeira_marca or ultima_qualquer


def _fmt_data_hora(dt):
    """datetime -> 'dd/mm/aaaa às HH:MM'."""
    return dt.strftime("%d/%m/%Y às %H:%M")


def _is_concluido(status):
    s = (status or "").lower()
    return any(x in s for x in STATUS_CONCLUIDO)


def _atividade_concluida(status_os):
    return (status_os or "").strip().lower() in STATUS_OS_CONCLUIDO


def _status_exibicao(status_raw):
    """Traduz o status bruto (statusOS do Fracttal / status de Falha) para
    o rótulo mostrado no relatório do cliente — evita imprimir jargão
    interno como 'Finalizada'/'Em Revisão' sem contexto. Corrigido
    22/07/2026 (ver ajuste manual do Fred no relatório Alves Lima Semana 30)."""
    s = (status_raw or "").strip().lower()
    if s == "finalizada":
        return "Concluída"
    if s in ("em revisão", "em revisao"):
        return "Em verificação"
    if s in STATUS_CONCLUIDO:
        return "Concluído"
    return (status_raw or "Em aberto").strip()


def _ticket_valido(ticket):
    """Só considera 'chamado real' se houver um número/identificador de verdade."""
    t = (ticket or "").strip().lower()
    return t not in ("", "n/a", "na", "-", "--", "nao", "não", "sem chamado", "sem ticket", "s/n")


def _ordenar_categorias(chaves):
    """Ordena rótulos de categoria pela ordem definida em ORDEM_CATEGORIAS_GERAIS;
    o que não estiver na lista vai para o final, em ordem alfabética."""
    def chave_ordenacao(rotulo):
        if rotulo in ORDEM_CATEGORIAS_GERAIS:
            return (0, ORDEM_CATEGORIAS_GERAIS.index(rotulo))
        return (1, rotulo)
    return sorted(chaves, key=chave_ordenacao)


# ── Coleta e agrupamento — Painel de Falhas (ocorrências) ──────────────────
#
# Toda ocorrência/atividade, seja qual for a origem, é normalizada para o
# mesmo formato de dict antes de entrar em "grupos":
#   {"usina", "equipamento", "falha", "status", "dt_abertura",
#    "dt_fechamento", "concluida", "origem"}
# Isso permite juntar Painel de Falhas + Painel de Atividades num único
# relatório sem duplicar as funções de geração de texto/slide.

def coletar_ocorrencias_semana(todos_valores, cliente, data_inicio, data_fim):
    """
    todos_valores: retorno de ws.get_all_values() do Painel de Falhas (linha 0 = cabeçalho)
    cliente: nome do cliente (comparação por 'contains', case-insensitive)
    data_inicio / data_fim: datetime (00:00 do dia inicial até 23:59 do dia final)

    CORRIGIDO 22/07/2026: só entra o que abriu OU fechou dentro do período
    selecionado (mesma correção aplicada em coletar_atividades_semana) —
    uma ocorrência aberta há muito tempo não deve reaparecer em todo
    relatório seguinte só por continuar sem status de concluído.

    Retorna dict: {"categoria": {"concluidas": [dict,...], "abertas": [dict,...]}, ...}
    """
    cliente_norm = _norm(cliente)
    grupos = {}

    for row in todos_valores[1:]:
        if len(row) <= COL_DATA_FECHAMENTO:
            row = row + [""] * (COL_DATA_FECHAMENTO + 1 - len(row))

        if cliente_norm not in _norm(row[COL_CLIENTE]):
            continue

        if _norm(row[COL_STATUS]) in STATUS_FALHA_EXCLUIR_RELATORIO:
            # já reportada como Atividade (com número de OS) — evita duplicar
            continue

        dt_abertura = _parse_data(row[COL_DATA_ABERTURA])
        concluida = _is_concluido(row[COL_STATUS])
        if concluida:
            # mesmo raciocínio de coletar_atividades_semana: prioriza a
            # última data real do Histórico sobre a Data de Fechamento
            # (que pode ter sido gravada só quando alguém confirmou o
            # fechamento no painel, não quando o problema foi resolvido
            # de fato).
            dt_fecham = _ultima_data_historico(row[COL_HISTORICO]) or _parse_data(row[COL_DATA_FECHAMENTO])
        else:
            dt_fecham = None

        fechou_na_semana = dt_fecham and data_inicio <= dt_fecham <= data_fim
        abriu_na_semana = dt_abertura and data_inicio <= dt_abertura <= data_fim

        if not (fechou_na_semana or abriu_na_semana):
            continue

        texto_categoria = f"{row[COL_EQUIP]} {row[COL_FALHA]}"
        categoria = _rotulo_categoria(texto_categoria)
        d = {
            "usina": row[COL_USINA].strip() or "Usina não informada",
            "equipamento": row[COL_EQUIP].strip() or "Equipamento não informado",
            "numero_os": row[COL_OS].strip() or None,
            "falha": row[COL_FALHA].strip() or "Falha não especificada",
            "status": row[COL_STATUS].strip() or "Em aberto",
            "dt_abertura": dt_abertura,
            "dt_fechamento": dt_fecham,
            "concluida": concluida,
            "origem": "falha",
        }
        grupos.setdefault(categoria, {"concluidas": [], "abertas": []})

        if fechou_na_semana or (concluida and abriu_na_semana):
            # fechou dentro da semana, OU abriu e já está concluída mas sem
            # Data de Fechamento gravada corretamente (mesmo fallback usado
            # em coletar_atividades_semana).
            grupos[categoria]["concluidas"].append(d)
        else:
            # abriu dentro da semana e ainda está em aberto — mostra como
            # aberta só nesta janela (não perpetuamente).
            grupos[categoria]["abertas"].append(d)

    return {k: v for k, v in grupos.items() if v["concluidas"] or v["abertas"]}


# ── Coleta e agrupamento — Painel de Atividades (OSs) ───────────────────────

def coletar_atividades_semana(todos_valores, cliente, data_inicio, data_fim):
    """
    todos_valores: retorno de ws.get_all_values() do Painel de Atividades
    (linha 0 = cabeçalho). A data de abertura é a Data de Criação da
    OS/atividade e a de fechamento é a Data de Conclusão; "concluída" é
    definido pelo statusOS (só "Finalizada" — ver STATUS_OS_CONCLUIDO).

    CORRIGIDO 22/07/2026: o relatório semanal só deve trazer o que
    realmente aconteceu (abriu OU fechou) DENTRO do período selecionado —
    uma atividade ainda aberta (ex.: presa em "Em Revisão" há semanas,
    criada bem antes do período) não deve reaparecer em todo relatório
    seguinte só por continuar sem "Finalizada". O caso relatado pelo Fred:
    OS 9173 (ABC Morada Nova), criada em 08/07 e parada em "Em Revisão" a
    100% desde então, aparecia num relatório de 16/07 a 22/07 mesmo sem
    nenhum evento dela ter ocorrido nessa janela.

    Retorna dict no mesmo formato de coletar_ocorrencias_semana, pronto
    para ser mesclado com mesclar_grupos().
    """
    cliente_norm = _norm(cliente)
    grupos = {}
    min_cols = ATIV_COL_STATUSOS + 1

    for row in todos_valores[1:]:
        if len(row) <= min_cols:
            row = row + [""] * (min_cols + 1 - len(row))

        if cliente_norm not in _norm(row[ATIV_COL_CLIENTE]):
            continue

        dt_abertura = _parse_data(row[ATIV_COL_DATA_CRIACAO])
        concluida = _atividade_concluida(row[ATIV_COL_STATUSOS])
        if concluida:
            # Prioriza a data real do Histórico sobre a Data de Conclusão
            # carimbada pela auditoria (que reflete quando NÓS detectamos o
            # "Finalizada", não quando aconteceu de fato — ver
            # _ultima_data_historico). Só cai para a Data de Conclusão se
            # não houver nada de útil no Histórico. Só faz sentido calcular
            # isso quando a OS está de fato concluída — senão qualquer
            # atualização recente de uma OS ainda em andamento poderia ser
            # confundida com um "fechamento" (bug evitado aqui de propósito).
            dt_fecham = _ultima_data_historico(row[ATIV_COL_HISTORICO]) or _parse_data(row[ATIV_COL_DATA_CONCLUSAO])
        else:
            dt_fecham = None

        fechou_na_semana = dt_fecham and data_inicio <= dt_fecham <= data_fim
        abriu_na_semana = dt_abertura and data_inicio <= dt_abertura <= data_fim

        if not (fechou_na_semana or abriu_na_semana):
            continue

        # categoriza pela descrição (onde aparecem palavras como "desligada",
        # "offline", "religamento" etc.) — o código do equipamento sozinho
        # (ex.: "2C-APG100") nunca bate em nenhum padrão, então cai sempre
        # em OUTRAS OCORRÊNCIAS e perde a categorização real (bug corrigido
        # em 13/07/2026). Usa equipamento+descrição juntos pra manter
        # também os casos que só têm palavra-chave no nome do equipamento.
        texto_categoria = f"{row[ATIV_COL_EQUIP]} {row[ATIV_COL_DESCRICAO]}"
        categoria = _rotulo_categoria(texto_categoria)
        status_display = row[ATIV_COL_STATUSOS].strip() or row[ATIV_COL_STATUS].strip() or "Em aberto"
        d = {
            "usina": row[ATIV_COL_USINA].strip() or "Usina não informada",
            "equipamento": row[ATIV_COL_EQUIP].strip() or "Equipamento não informado",
            "numero_os": row[ATIV_COL_NUMEROOS].strip() or None,
            "falha": row[ATIV_COL_DESCRICAO].strip() or "Sem descrição",
            "status": status_display,
            "dt_abertura": dt_abertura,
            "dt_fechamento": dt_fecham,
            "concluida": concluida,
            "origem": "atividade",
        }
        grupos.setdefault(categoria, {"concluidas": [], "abertas": []})

        if fechou_na_semana or (concluida and abriu_na_semana):
            # fechou dentro da semana, OU abriu E já está "Finalizada" mas
            # a Data de Conclusão ficou vazia/fora da janela (bug histórico
            # corrigido em 13/07/2026 — atividades concluídas não gravavam
            # essa data). Sem este fallback, a atividade some do relatório
            # sem deixar rastro.
            grupos[categoria]["concluidas"].append(d)
        else:
            # abriu dentro da semana mas ainda não fechou — mostra como
            # aberta/em andamento, só nesta janela (não perpetuamente).
            grupos[categoria]["abertas"].append(d)

    return {k: v for k, v in grupos.items() if v["concluidas"] or v["abertas"]}


# ── NOVO PADRÃO DE RELATÓRIO (confirmado com Fred 23/07/2026) ──────────────
#
# Agrupamento por USINA (não por categoria de equipamento). Só entram OSs
# vinculadas à Fracttal (numeroOS preenchido) cujo statusOS esteja em
# "Em Revisão" (= aba "Em Verificação" do painel) ou "Finalizada"
# (= aba "Concluídas") — nunca "Em Processo"/"Cancelada". Ambas aparecem
# no relatório com o rótulo "Concluída" (mesmo padrão do relatório de
# referência corrigido manualmente pelo Fred). A data usada para decidir
# se a OS pertence à semana do relatório é a mesma já extraída do
# Histórico (marca de "100%"/mudança de status), com fallback pra Data de
# Conclusão — decisão explícita do Fred (23/07/2026): usar só dado que já
# está salvo na planilha, sem consulta ao vivo à Fracttal por OS.
# ── Classificação de desligamento — PORTADA de isDesligamentoAtividade()
#    (index.html do painel, linhas ~2709-2747). Usa os mesmos padrões e a
#    mesma lógica de negação (_clauseTemNegacao/_matchSemNegacao) do
#    dashboard, pra este relatório classificar desligamento exatamente
#    como o painel classifica — a tentativa anterior reaproveitava
#    _rotulo_categoria (feita pra CATEGORIZAR EQUIPAMENTO, não pra
#    detectar desligamento) e deixava passar casos reais como
#    "Religamento de usina" (23/07/2026: corrigido após comparar a saída
#    real do endpoint com o relatório de referência do Fred — OS 9701,
#    9750, 9666, 9591, 9396 apareciam em ATIVIDADES DA SEMANA quando
#    deveriam estar em DESLIGAMENTOS).

def _clausula_tem_negacao(texto, indice_match):
    inicio = indice_match
    while inicio > 0 and texto[inicio - 1] not in ".!?\n":
        inicio -= 1
    clausula = texto[inicio:indice_match]
    return bool(re.search(r"\b(não|nao|sem|dispensad[ao]|descartad[ao]|nunca)\b", clausula, re.IGNORECASE))


def _match_sem_negacao(padrao, texto):
    for m in re.finditer(padrao, texto, re.IGNORECASE):
        if not _clausula_tem_negacao(texto, m.start()):
            return True
    return False


_PADROES_DESLIGAMENTO_ATIVIDADE = [
    r"(?:usina|ufv)\s+(?:\w+\s+){0,3}(?:desligad[ao]|parad[ao]|sem\s+energia|desenergizad[ao]|offline)",
    r"(?:desligad[ao]|parad[ao]|offline)\s+(?:\w+\s+){0,3}(?:usina|ufv)",
    r"sem\s+possibilidade\s+de\s+(?:monitoramento|religamento)",
    r"(?:usina|ufv)\s+(?:\w+\s+){0,3}sem\s+comunica[cç][ãa]o",
    r"sem\s+comunica[cç][ãa]o\s+(?:\w+\s+){0,3}(?:usina|ufv)",
    r"desligamento\s+(?:total\s+)?(?:da|de)\s+(?:usina|ufv)",
    r"religamento\s+(?:da|de)\s+(?:usina|ufv)",
    r"transformador\s+(?:\w+\s+){0,3}(?:da|de)\s+(?:usina|ufv)\s+(?:\w+\s+){0,3}(?:desligad[ao]|parad[ao])",
    r"trafo\s+(?:\w+\s+){0,3}(?:da|de)\s+(?:usina|ufv)\s+(?:\w+\s+){0,3}(?:desligad[ao]|parad[ao])",
]


def _e_desligamento_atividade(descricao, equipamento):
    dc = f"{descricao} {equipamento}"
    return any(_match_sem_negacao(p, dc) for p in _PADROES_DESLIGAMENTO_ATIVIDADE)


STATUS_OS_ELEGIVEIS_RELATORIO = ["finalizada", "em revisão", "em revisao"]

VERDE_STATUS = "00B050"  # verde do relatório PPTX de cliente — não é o A1CA40 da marca


def coletar_atividades_e_desligamentos_por_usina(todos_valores, cliente, data_inicio, data_fim):
    """
    todos_valores: ws.get_all_values() do Painel de Atividades (linha 0 = cabeçalho).

    Retorna (atividades_por_usina, desligamentos_por_usina), cada um um
    dict {usina: [{"descricao":, "numero_os":}, ...]}. Itens cuja
    categoria (equipamento+descrição) bate no padrão de desligamento vão
    para o segundo dict em vez do primeiro — nunca duplicados.
    """
    cliente_norm = _norm(cliente)
    atividades, desligamentos = {}, {}
    min_cols = ATIV_COL_STATUSOS + 1

    for row in todos_valores[1:]:
        if len(row) <= min_cols:
            row = row + [""] * (min_cols + 1 - len(row))

        # Campo Cliente da própria OS manda — nunca o nome da usina/site.
        # Isso já exclui automaticamente OS de outro cliente na mesma usina
        # física (ex.: OS da GreenYellow numa usina RENOGRID).
        if cliente_norm not in _norm(row[ATIV_COL_CLIENTE]):
            continue

        numero_os = row[ATIV_COL_NUMEROOS].strip()
        if not numero_os:
            continue  # relatório só traz OS vinculada à Fracttal

        status_os = row[ATIV_COL_STATUSOS].strip().lower()
        if status_os not in STATUS_OS_ELEGIVEIS_RELATORIO:
            continue

        dt_marco = _ultima_data_historico(row[ATIV_COL_HISTORICO]) or _parse_data(row[ATIV_COL_DATA_CONCLUSAO])
        if not dt_marco or not (data_inicio <= dt_marco <= data_fim):
            continue

        usina = row[ATIV_COL_USINA].strip() or "Usina não informada"
        descricao = row[ATIV_COL_DESCRICAO].strip() or "Sem descrição"
        equipamento = row[ATIV_COL_EQUIP].strip()
        if len(descricao) > 140:
            descricao = descricao[:137].rstrip() + "..."

        item = {"descricao": descricao, "numero_os": numero_os}

        if _e_desligamento_atividade(descricao, equipamento):
            desligamentos.setdefault(usina, []).append(item)
        else:
            atividades.setdefault(usina, []).append(item)

    return atividades, desligamentos


def _formatar_item_atividade(it):
    return [
        {"texto": f'{it["descricao"]} – ', "bold": False},
        {"texto": f'OS {it["numero_os"]}', "bold": True},
        {"texto": " – ", "bold": False},
        {"texto": "Concluída", "bold": False, "color": VERDE_STATUS},
    ]


def _formatar_item_desligamento(it):
    return [
        {"texto": "Desligamento - ", "bold": False},
        {"texto": f'OS {it["numero_os"]}', "bold": True},
        {"texto": " – ", "bold": False},
        {"texto": "Concluída.", "bold": False, "color": VERDE_STATUS},
    ]


# ── Formatação explícita (Poppins 20pt, numeração automática) ──────────────
# Não depende de herdar a formatação do modelo (o .pptx-base usa marcador
# "Ø" avulso, sem numeração) — reproduz explicitamente o padrão exato do
# relatório de referência corrigido manualmente pelo Fred (23/07/2026):
# título numerado (arábico), usina numerada (romano maiúsculo), item com
# marcador "•", tudo Poppins 20pt, 6pt de espaço antes de cada parágrafo.

def _pPr_xml(nivel, start_at=None, bullet_char="•"):
    ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    spc = '<a:spcBef><a:spcPts val="600"/></a:spcBef>'
    sa = f' startAt="{start_at}"' if start_at else ""
    if nivel == "titulo":
        corpo = f'<a:buFont typeface="+mj-lt"/><a:buAutoNum type="arabicPeriod"{sa}/>'
        attrs = 'marL="756920" marR="347980" indent="-457200" algn="l"'
    elif nivel == "usina":
        corpo = f'<a:buFont typeface="+mj-lt"/><a:buAutoNum type="romanUcPeriod"{sa}/>'
        attrs = 'marL="814070" marR="347980" indent="-514350" algn="l"'
    elif nivel == "item":
        fonte_bullet = "Arial" if bullet_char == "•" else "+mn-lt"
        corpo = f'<a:buFont typeface="{fonte_bullet}" pitchFamily="34" charset="0"/><a:buChar char="{bullet_char}"/>'
        attrs = 'marL="642620" marR="347980" indent="-342900" algn="l"'
    elif nivel == "subtitulo":
        corpo = "<a:buNone/>"
        attrs = 'marL="299720" marR="347980" algn="l"'
    else:  # blank — linha em branco entre blocos de usina
        corpo = ""
        attrs = 'marL="299720" marR="347980" algn="l"'
    return parse_xml(f"<a:pPr {ns} {attrs}>{spc}{corpo}</a:pPr>")


def _add_paragrafo(tf, nivel, runs, first=False, start_at=None, bullet_char="•"):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p_el = p._p
    velho_pPr = p_el.find(qn("a:pPr"))
    if velho_pPr is not None:
        p_el.remove(velho_pPr)
    p_el.insert(0, _pPr_xml(nivel, start_at=start_at, bullet_char=bullet_char))
    for r in runs:
        run = p.add_run()
        run.text = r.get("texto", "")
        run.font.name = "Poppins"
        run.font.size = Pt(20)
        run.font.bold = bool(r.get("bold", False))
        if r.get("italic"):
            run.font.italic = True
        if r.get("color"):
            run.font.color.rgb = RGBColor.from_string(r["color"])
    return p


def _gerar_blocos_usina(usinas_ordenadas, dados_por_usina, formatar_item, texto_vazio, bullet_char):
    blocos = []
    for usina in usinas_ordenadas:
        itens = dados_por_usina.get(usina, [])
        paragrafos = [{"nivel": "usina", "runs": [{"texto": usina, "bold": True}]}]
        if itens:
            for it in itens:
                paragrafos.append({"nivel": "item", "runs": formatar_item(it), "bullet_char": bullet_char})
        else:
            paragrafos.append({"nivel": "item",
                                "runs": [{"texto": texto_vazio, "italic": True}],
                                "bullet_char": bullet_char})
        blocos.append({"usina": usina, "paragrafos": paragrafos, "linhas": 2 + len(itens)})
    return blocos


def _renderizar_topico_usinas(prs, numero_topico, titulo_base, usinas_ordenadas, dados_por_usina,
                               formatar_item, texto_vazio, bullet_char, max_linhas):
    blocos = _gerar_blocos_usina(usinas_ordenadas, dados_por_usina, formatar_item, texto_vazio, bullet_char)

    paginas, atual, linhas_atual = [], [], 0
    for b in blocos:
        if atual and linhas_atual + b["linhas"] > max_linhas:
            paginas.append(atual)
            atual, linhas_atual = [], 0
        atual.append(b)
        linhas_atual += b["linhas"]
    if atual:
        paginas.append(atual)
    if not paginas:
        paginas = [[]]

    numero_romano = 1
    for i, pagina_blocos in enumerate(paginas):
        titulo = titulo_base if i == 0 else f"{titulo_base} – CONTINUAÇÃO"
        novo = _duplicate_slide(prs, 2)
        _remover_fotos(novo)
        shp_titulo_secao = _find_shape(novo, "DESLIGAMENTOS")
        shp_corpo = _find_shape(novo, "Foram registradas")
        if shp_titulo_secao:
            tf_t = shp_titulo_secao.text_frame
            tf_t.clear()
            _add_paragrafo(tf_t, "titulo", [{"texto": titulo, "bold": True}], first=True, start_at=numero_topico)
        if shp_corpo:
            tf = shp_corpo.text_frame
            tf.clear()
            primeiro = True
            for b in pagina_blocos:
                for par in b["paragrafos"]:
                    if par["nivel"] == "usina":
                        _add_paragrafo(tf, "usina", par["runs"], first=primeiro, start_at=numero_romano)
                        numero_romano += 1
                    else:
                        _add_paragrafo(tf, "item", par["runs"], first=primeiro,
                                        bullet_char=par.get("bullet_char", "•"))
                    primeiro = False
                _add_paragrafo(tf, "blank", [{"texto": ""}], first=False)
    return novo


PAUTAS_GERAIS_FIXAS = ["ATIVIDADES DA SEMANA", "DESLIGAMENTOS", "CHAMADOS E PROTOCOLOS",
                       "OUTRAS ATIVIDADES", "ZELADORIA"]


def _renderizar_pautas_gerais(prs):
    """Slide 2 — pauta fixa e universal (vale para todos os clientes,
    confirmado com Fred 23/07/2026), substitui a antiga 'Ata da reunião'
    com pauta variável por cliente."""
    ata = _duplicate_slide(prs, 1)
    shp_titulo = _find_shape(ata, "Ata da reuni")
    if shp_titulo:
        _set_text_preservando_estilo(shp_titulo, "Gestão e Operação de Ativos\nPautas Gerais")
    shp_corpo = _find_shape(ata, "Desligamentos")
    if shp_corpo:
        tf = shp_corpo.text_frame
        tf.clear()
        for i, topico in enumerate(PAUTAS_GERAIS_FIXAS):
            _add_paragrafo(tf, "titulo", [{"texto": f"{topico};", "bold": False}], first=(i == 0), start_at=i + 1)
    return ata


def _renderizar_secao_placeholder(prs, numero_topico, titulo, corpo_callback=None):
    """Slide só com o título numerado — usado em CHAMADOS E PROTOCOLOS,
    OUTRAS ATIVIDADES e ZELADORIA, onde o conteúdo é preenchido manualmente
    pelo Fred (não vem do dashboard)."""
    novo = _duplicate_slide(prs, 2)
    _remover_fotos(novo)
    shp_titulo_secao = _find_shape(novo, "DESLIGAMENTOS")
    shp_corpo = _find_shape(novo, "Foram registradas")
    if shp_titulo_secao:
        tf_t = shp_titulo_secao.text_frame
        tf_t.clear()
        _add_paragrafo(tf_t, "titulo", [{"texto": titulo, "bold": True}], first=True, start_at=numero_topico)
    if shp_corpo:
        tf = shp_corpo.text_frame
        tf.clear()
        if corpo_callback:
            corpo_callback(tf)
    return novo


def mesclar_grupos(*grupos_varios):
    """Mescla dois ou mais dicts no formato de coletar_ocorrencias_semana
    (ex.: Painel de Falhas + Painel de Atividades) em um único dict de
    categorias, somando as listas de concluídas/abertas de cada fonte."""
    resultado = {}
    for grupos in grupos_varios:
        for cat, dados in (grupos or {}).items():
            alvo = resultado.setdefault(cat, {"concluidas": [], "abertas": []})
            alvo["concluidas"].extend(dados.get("concluidas", []))
            alvo["abertas"].extend(dados.get("abertas", []))
    return resultado


def listar_usinas_cliente(todos_valores, cliente):
    """Lista todas as usinas do cliente (histórico completo, não só a semana)."""
    cliente_norm = _norm(cliente)
    usinas = []
    vistos = set()
    for row in todos_valores[1:]:
        if len(row) <= COL_USINA:
            continue
        if cliente_norm not in _norm(row[COL_CLIENTE]):
            continue
        usina = row[COL_USINA].strip()
        chave = _norm(usina)
        if usina and chave not in vistos:
            vistos.add(chave)
            usinas.append(usina)
    return sorted(usinas)


def coletar_chamados_abertos(todos_valores, cliente):
    """Ocorrências do cliente (Painel de Falhas) com número de chamado/ticket
    fabricante válido e ainda não concluídas."""
    cliente_norm = _norm(cliente)
    chamados = []
    for row in todos_valores[1:]:
        if len(row) <= COL_DATA_FECHAMENTO:
            row = row + [""] * (COL_DATA_FECHAMENTO + 1 - len(row))
        if cliente_norm not in _norm(row[COL_CLIENTE]):
            continue
        ticket = row[COL_TICKET].strip()
        if not _ticket_valido(ticket) or _is_concluido(row[COL_STATUS]):
            continue
        chamados.append(row)
    return chamados


# ── Geração de conteúdo (offline, sem IA) — parágrafos com negrito por usina ─

def _linha_ocorrencia(d):
    """Linha resumida: Descrição – OS nº (quando houver). Status: X.
    (sem despejar causa/ação/histórico). Corrigido 22/07/2026: usar o
    número da OS (legível pro cliente) em vez do código interno do
    equipamento sempre que disponível, e traduzir o status bruto."""
    falha = d["falha"]
    if len(falha) > 180:
        falha = falha[:177].rstrip() + "..."
    status = _status_exibicao(d["status"])
    numero_os = d.get("numero_os")
    if numero_os:
        return f"{falha} – OS {numero_os}. Status: {status}."
    equip = d["equipamento"]
    return f"{equip} – {falha}. Status: {status}."


def gerar_paragrafos_categoria(categoria, dados):
    """
    Agrupa as ocorrências/atividades de uma categoria por usina (nome da
    usina em negrito, seguido de um bullet resumido por ocorrência).
    Retorna lista de dicts: {"texto": str, "bold": bool}
    """
    todas = list(dados["concluidas"]) + list(dados["abertas"])
    por_usina = {}
    for d in todas:
        por_usina.setdefault(d["usina"], []).append(d)

    paragrafos = []
    for usina in sorted(por_usina.keys()):
        paragrafos.append({"texto": usina, "bold": True})
        for d in por_usina[usina]:
            paragrafos.append({"texto": "• " + _linha_ocorrencia(d), "bold": False})
    return paragrafos


# ── Desligamentos: sempre com data/hora real (aberto via WhatsApp já traz a
# hora do desligamento; ao normalizar, a hora de conclusão) ─────────────────

def _linha_desligamento_runs(d):
    """
    Runs (mistura de negrito/normal na mesma linha) para uma ocorrência de
    desligamento, incluindo data/hora real de abertura e, se concluída, de
    conclusão. Só cai no aviso em negrito se a data realmente estiver
    ausente na origem (caso excepcional — normalmente já vem preenchida
    pela automação do WhatsApp, ou pela Data de Criação/Conclusão da OS).
    """
    equip = d["equipamento"]
    falha = d["falha"]
    status = _status_exibicao(d["status"])
    concluida = d["concluida"]
    dt_abertura = d["dt_abertura"]
    dt_fecham = d["dt_fechamento"]

    runs = [{"texto": f"• {equip} – {falha}. Status: {status}. ", "bold": False}]

    if dt_abertura:
        runs.append({"texto": f"Desligada em {_fmt_data_hora(dt_abertura)}.", "bold": False})
    else:
        runs.append({"texto": "ACRESCENTAR A DATA E HORA DO DESLIGAMENTO", "bold": True})

    if concluida:
        if dt_fecham:
            runs.append({"texto": f" Religada/concluída em {_fmt_data_hora(dt_fecham)}.", "bold": False})
        else:
            runs.append({"texto": " ACRESCENTAR A DATA E HORA DA CONCLUSÃO", "bold": True})

    return runs


def gerar_paragrafos_desligamentos(dados):
    """Igual a gerar_paragrafos_categoria, mas com data/hora real por ocorrência."""
    todas = list(dados["concluidas"]) + list(dados["abertas"])
    por_usina = {}
    for d in todas:
        por_usina.setdefault(d["usina"], []).append(d)

    paragrafos = []
    for usina in sorted(por_usina.keys()):
        paragrafos.append({"runs": [{"texto": usina, "bold": True}]})
        for d in por_usina[usina]:
            paragrafos.append({"runs": _linha_desligamento_runs(d)})
    return paragrafos


# ── Chamados em aberto: uma linha por chamado, usina em negrito inline ─────

def gerar_paragrafos_chamados_fabricante(chamados):
    """Uma linha por chamado da aba ChamadosFabricante (dashboard novo,
    17/07/2026) — mostra Ativo, Identificação do Equipamento, Ticket/RMA,
    Status e a Nota do dashboard (quando o Fred escreveu alguma no popup
    de detalhe). Recebe a lista JÁ filtrada por cliente e sem os
    concluídos — este código só formata, não filtra."""
    if not chamados:
        return [{"texto": "Nenhum chamado em aberto nesta semana.", "bold": False}]

    por_usina = {}
    for c in chamados:
        usina = (c.get("UFV") or "Usina não informada").strip()
        por_usina.setdefault(usina, []).append(c)

    paragrafos = []
    for usina in sorted(por_usina.keys()):
        paragrafos.append({"runs": [{"texto": usina, "bold": True}]})
        for c in por_usina[usina]:
            ativo = (c.get("Ativo") or "").strip()
            equip = (c.get("Identificação do Equipamento") or "").strip()
            ticket = (c.get("Ticket/RMA") or "").strip()
            status = (c.get("Status") or "Em análise").strip()
            nota = (c.get("NotaDashboard") or "").strip()

            descricao_equip = equip or ativo or "Equipamento não informado"
            if equip and ativo and ativo.lower() not in equip.lower():
                descricao_equip = f"{ativo} — {equip}"

            texto_principal = f"- {descricao_equip} – Ticket {ticket or 's/n'} – {status}."
            runs = [{"texto": texto_principal, "bold": False}]
            if nota:
                runs.append({"texto": f" Obs.: {nota}", "bold": False})
            paragrafos.append({"runs": runs})
    return paragrafos


def gerar_paragrafos_chamados(chamados):
    """Uma linha por chamado válido: '- Usina – Case #ticket -> resumo -> status.'
    (chamados vêm sempre do Painel de Falhas — Painel de Atividades não tem
    número de chamado/ticket de fabricante, só número de OS interno)."""
    if not chamados:
        return [{"texto": "Nenhum chamado em aberto nesta semana.", "bold": False}]

    por_usina = {}
    for row in chamados:
        usina = row[COL_USINA].strip() or "Usina não informada"
        por_usina.setdefault(usina, []).append(row)

    paragrafos = []
    for usina in sorted(por_usina.keys()):
        for row in por_usina[usina]:
            ticket = row[COL_TICKET].strip()
            resumo = row[COL_FALHA].strip() or row[COL_CAUSA].strip() or "Sem descrição detalhada"
            status = row[COL_STATUS].strip() or "Em análise"
            runs = [
                {"texto": "- ", "bold": False},
                {"texto": usina, "bold": True},
                {"texto": f" – Case #{ticket} -> {resumo} -> {status}.", "bold": False},
            ]
            paragrafos.append({"runs": runs})
    return paragrafos


# Layout real da aba "Zeladoria" (gid 987654321): linha 1 = título dos grupos,
# linha 2 = subcabeçalhos, dados a partir da linha 3.
# Colunas (0-based): A Cliente, B Usina, depois 4 grupos de 4 colunas cada
# (Última Data, Próxima Data, Quantidade, Status): Roçada, Poda Química, Lavagem dos Módulos, Controle de Pragas.
ZEL_COL_CLIENTE, ZEL_COL_USINA = 0, 1
ZEL_GRUPOS = [
    ("Roçada",               2),
    ("Poda Química",         6),
    ("Lavagem dos Módulos", 10),
    ("Controle de Pragas",  14),
]


# Clientes onde uma célula de Zeladoria vazia deve exibir "Acompanhamento"
# em vez de "Sem informação" — combinado com Fred em 22/07/2026: ele ajusta
# manualmente quando necessário, mas o padrão nesses dois clientes deve
# partir de "Acompanhamento".
CLIENTES_ZELADORIA_DEFAULT_ACOMPANHAMENTO = ["gd energy", "alves lima"]


def coletar_zeladoria(zeladoria_valores, cliente):
    """
    zeladoria_valores: ws.get_all_values() da aba Zeladoria (linhas 1-2 = cabeçalho).
    Retorna lista de dicts: [{"usina": ..., "grupos": [{"nome":, "status":, "ultima_data":}, ...]}, ...]
    """
    cliente_norm = _norm(cliente)
    usa_default_acompanhamento = any(c in cliente_norm for c in CLIENTES_ZELADORIA_DEFAULT_ACOMPANHAMENTO)
    resultado = []
    for row in zeladoria_valores[2:]:
        if len(row) <= ZEL_GRUPOS[-1][1] + 3:
            row = row + [""] * (ZEL_GRUPOS[-1][1] + 4 - len(row))
        if not row[ZEL_COL_USINA].strip():
            continue
        if cliente_norm not in _norm(row[ZEL_COL_CLIENTE]):
            continue

        grupos_usina = []
        for nome, col_ini in ZEL_GRUPOS:
            ultima_data = row[col_ini].strip()
            status = row[col_ini + 3].strip()
            # A célula pode estar vazia OU conter literalmente o texto
            # "Sem informação(ões)" já gravado como dado na planilha (por
            # uma rodada anterior do relatório, ou por preenchimento
            # manual) — nesse caso "if not status" nunca disparava, porque
            # a célula não está tecnicamente vazia. Tratamos esse texto
            # como equivalente a vazio pra aplicar o padrão "Acompanhamento"
            # do mesmo jeito. Corrigido 22/07/2026.
            if _norm(status) in ("", "sem informacao", "sem informacoes") and usa_default_acompanhamento:
                status = "Acompanhamento"
            grupos_usina.append({"nome": nome, "status": status, "ultima_data": ultima_data})

        resultado.append({"usina": row[ZEL_COL_USINA].strip(), "grupos": grupos_usina})
    return resultado


def _status_zeladoria_para_tabela(item, nome_grupo):
    """Texto de uma célula da tabela de Zeladoria (só o status, conforme
    pedido — sem despejar a última data dentro da célula para não poluir)."""
    grupo = next((g for g in item["grupos"] if g["nome"] == nome_grupo), None)
    if not grupo:
        return "Sem informação"
    return grupo["status"] or "Sem informação"


def paginar_zeladoria(zeladoria_usinas, max_por_pagina=ZELADORIA_MAX_USINAS_POR_PAGINA):
    """Divide as usinas da Zeladoria em páginas de tabela, pra nunca estourar
    o slide (era o problema antigo: texto cortado com muitas usinas)."""
    if not zeladoria_usinas:
        return [[]]
    return [zeladoria_usinas[i:i + max_por_pagina] for i in range(0, len(zeladoria_usinas), max_por_pagina)]


# ── Clonagem de slides (mantém 100% do estilo do modelo Grid Co.) ──────────

def _duplicate_slide(prs, index):
    """Duplica um slide existente do modelo (mesmo layout, cores, fontes, formas)."""
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)

    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    id_map = {}
    for rId, rel in source.part.rels.items():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        new_rid = dest.part.rels._add_relationship(rel.reltype, rel._target, rel.is_external)
        id_map[rId] = new_rid

    for shp in source.shapes:
        newel = copy.deepcopy(shp._element)
        for el in newel.iter():
            for attr in ("embed", "link", "id"):
                full = qn(f"r:{attr}")
                val = el.get(full)
                if val and val in id_map:
                    el.set(full, id_map[val])
        dest.shapes._spTree.append(newel)
    return dest


def _extrair_formato_base(shape):
    """Extrai fonte/tamanho/itálico/cor do 1º run existente na shape (pra reaplicar depois)."""
    tf = shape.text_frame
    primeiro_par = tf.paragraphs[0]
    if not primeiro_par.runs:
        return None
    font = primeiro_par.runs[0].font
    cor_rgb = None
    try:
        if font.color and font.color.type is not None:
            cor_rgb = font.color.rgb  # levanta AttributeError se for cor de tema (schemeClr)
    except AttributeError:
        cor_rgb = None
    return (font.name, font.size, font.italic, cor_rgb)


def _extrair_ppr_bullet(shape):
    """
    Extrai o XML de formatação de parágrafo (<a:pPr>, que carrega a definição
    do bullet/marcador ">" do modelo) do 1º parágrafo da shape, ANTES de
    qualquer tf.clear(). Sem isso, só a 1ª linha herdava o bullet do modelo —
    as linhas adicionadas via add_paragraph() vinham sem marcador nenhum.
    """
    tf = shape.text_frame
    p_origem = tf.paragraphs[0]._p
    pPr = p_origem.find(qn("a:pPr"))
    return copy.deepcopy(pPr) if pPr is not None else None


def _aplicar_ppr_bullet(paragrafo, ppr_modelo):
    """Aplica (substituindo) o <a:pPr> de ppr_modelo num parágrafo já criado."""
    if ppr_modelo is None:
        return
    p_destino = paragrafo._p
    pPr_atual = p_destino.find(qn("a:pPr"))
    novo = copy.deepcopy(ppr_modelo)
    if pPr_atual is not None:
        p_destino.remove(pPr_atual)
    p_destino.insert(0, novo)


def _set_text_preservando_estilo(shape, novo_texto):
    """Substitui o texto de um text box mantendo a formatação (incl. negrito) do 1º run
    E o marcador de bullet (">") do modelo em TODAS as linhas, não só na 1ª."""
    tf = shape.text_frame
    primeiro_par = tf.paragraphs[0]
    bold_original = primeiro_par.runs[0].font.bold if primeiro_par.runs else None
    fmt = _extrair_formato_base(shape)
    ppr_bullet = _extrair_ppr_bullet(shape)

    tf.clear()
    linhas = novo_texto.split("\n")
    for i, linha in enumerate(linhas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _aplicar_ppr_bullet(p, ppr_bullet)
        run = p.add_run()
        run.text = linha
        if fmt:
            if fmt[0]:
                run.font.name = fmt[0]
            if fmt[1]:
                run.font.size = fmt[1]
            if fmt[2] is not None:
                run.font.italic = fmt[2]
            if fmt[3]:
                run.font.color.rgb = fmt[3]
        if bold_original is not None:
            run.font.bold = bold_original


def _set_paragrafos_com_estilo(shape, paragrafos):
    """
    Substitui o conteúdo de um text box por uma lista de parágrafos, permitindo
    alternar negrito por parágrafo (usado para destacar nome de usina/categoria)
    OU, quando o item tiver a chave "runs", misturar negrito/normal dentro da
    MESMA linha (usado em Desligamentos e Chamados em Aberto).

    paragrafos: lista de dicts, em um dos dois formatos:
      - {"texto": str, "bold": bool}                      -> parágrafo simples
      - {"runs": [{"texto": str, "bold": bool}, ...]}      -> parágrafo com runs mistos
    """
    fmt = _extrair_formato_base(shape)
    ppr_bullet = _extrair_ppr_bullet(shape)
    tf = shape.text_frame
    tf.clear()

    if not paragrafos:
        paragrafos = [{"texto": "", "bold": False}]

    for i, p in enumerate(paragrafos):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _aplicar_ppr_bullet(par, ppr_bullet)
        runs_def = p.get("runs") or [{"texto": p.get("texto", ""), "bold": p.get("bold", False)}]
        for r in runs_def:
            run = par.add_run()
            run.text = r.get("texto", "")
            if fmt:
                if fmt[0]:
                    run.font.name = fmt[0]
                if fmt[1]:
                    run.font.size = fmt[1]
                if fmt[2] is not None:
                    run.font.italic = fmt[2]
                if fmt[3]:
                    run.font.color.rgb = fmt[3]
            run.font.bold = bool(r.get("bold", False))


def _formatar_celula_tabela(celula, fmt, negrito):
    """Aplica a fonte/tamanho extraídos do template (fmt) numa célula de tabela."""
    celula.text_frame.word_wrap = True
    for p in celula.text_frame.paragraphs:
        p.alignment = p.alignment  # no-op, mantém padrão do template de tabela
        if not p.runs:
            p.add_run()
        for r in p.runs:
            if fmt:
                if fmt[0]:
                    r.font.name = fmt[0]
                if fmt[1]:
                    r.font.size = fmt[1]
            r.font.bold = negrito


def _remover_fotos(slide):
    """Remove os placeholders de foto (o gerador não tem fotos de campo)."""
    for shp in list(slide.shapes):
        if shp.shape_type == 13 and (shp.name or "").lower().startswith("imagem"):
            shp._element.getparent().remove(shp._element)


def _find_shape(slide, contem_texto):
    for shp in slide.shapes:
        if shp.has_text_frame and contem_texto.lower() in shp.text_frame.text.lower():
            return shp
    return None


# ── Montagem do PPTX final ──────────────────────────────────────────────────

def _extrair_orcamento_linhas(shape, fmt):
    """
    Estima quantas linhas cabem na caixa do modelo e quantos caracteres cabem
    por linha, a partir da altura/largura reais da shape (EMU) e do tamanho
    de fonte do template — substitui o orçamento por "número de caracteres",
    que não refletia a quebra de linha real e deixava o conteúdo estourar
    a página em categorias com texto mais longo.
    """
    altura_pt = shape.height / 12700
    largura_pt = shape.width / 12700
    tamanho_fonte_pt = fmt[1].pt if fmt and fmt[1] else 18
    altura_linha_pt = tamanho_fonte_pt * 1.35  # inclui folga de espaçamento entre parágrafos
    largura_media_char_pt = tamanho_fonte_pt * 0.52  # estimativa conservadora (fonte não é monoespaçada)
    max_linhas = max(1, int(altura_pt // altura_linha_pt))
    chars_por_linha = max(15, int(largura_pt // largura_media_char_pt))
    max_linhas = max(1, int(max_linhas * 0.8))  # margem de segurança de 20%
    return max_linhas, chars_por_linha


def _texto_paragrafo(paragrafo):
    if "runs" in paragrafo:
        return "".join(r.get("texto", "") for r in paragrafo["runs"])
    return paragrafo.get("texto", "")


def _linhas_ocupadas(paragrafo, chars_por_linha):
    texto = _texto_paragrafo(paragrafo)
    if not texto:
        return 1
    return max(1, -(-len(texto) // chars_por_linha))  # divisão com arredondamento pra cima


def _linhas_totais(paragrafos, chars_por_linha):
    return sum(_linhas_ocupadas(p, chars_por_linha) for p in paragrafos)


def gerar_paragrafos_multi_categoria(dados_por_categoria):
    """Igual a gerar_paragrafos_categoria, mas mesclando várias categorias
    (ex.: Comunicações + Controle de Vegetação + Outras Ocorrências) numa
    única lista agrupada só por usina, sem subtítulo de categoria — usado
    quando a pauta do relatório não distingue esses sub-tópicos e tudo cai
    sob um único tópico "ATIVIDADES SEMANAIS" (combinado 22/07/2026: a
    pauta fixa não tem "Controle de Vegetação" nem "Outras Ocorrências"
    como itens separados, então o conteúdo não pode ficar em páginas com
    esses títulos)."""
    por_usina = {}
    for dados in dados_por_categoria.values():
        todas = list(dados["concluidas"]) + list(dados["abertas"])
        for d in todas:
            por_usina.setdefault(d["usina"], []).append(d)

    paragrafos = []
    for usina in sorted(por_usina.keys()):
        paragrafos.append({"texto": usina, "bold": True})
        for d in por_usina[usina]:
            paragrafos.append({"texto": "• " + _linha_ocorrencia(d), "bold": False})
    return paragrafos


def paginar_paragrafos_simples(paragrafos, max_linhas, chars_por_linha):
    """Pagina uma lista plana de parágrafos (sem cabeçalho de categoria) por
    orçamento de linhas — usado pela seção mesclada ATIVIDADES SEMANAIS."""
    paginas, atual, atual_linhas = [], [], 0
    for p in paragrafos:
        linhas = _linhas_ocupadas(p, chars_por_linha)
        if atual and (atual_linhas + linhas > max_linhas):
            paginas.append(atual)
            atual, atual_linhas = [], 0
        atual.append(p)
        atual_linhas += linhas
    if atual:
        paginas.append(atual)
    return paginas or [[]]


def agrupar_em_paginas(grupos, max_linhas, chars_por_linha, ordem=None):
    """
    Agrupa categorias em páginas por orçamento de LINHAS (estimadas a partir
    da caixa real do modelo). Retorna lista de páginas, cada página = lista
    de tuplas (categoria, paragrafos).
    `ordem`, se fornecida, define a sequência das categorias (categorias fora
    dela vão para o final, em ordem alfabética).
    """
    chaves = _ordenar_categorias(grupos.keys()) if ordem is None else ordem
    blocos = [(cat, gerar_paragrafos_categoria(cat, grupos[cat])) for cat in chaves if cat in grupos]
    paginas, atual, atual_linhas = [], [], 0
    for cat, paragrafos in blocos:
        linhas_bloco = 1 + _linhas_totais(paragrafos, chars_por_linha)  # +1 = cabeçalho da categoria
        if atual and (atual_linhas + linhas_bloco > max_linhas):
            paginas.append(atual)
            atual, atual_linhas = [], 0
        atual.append((cat, paragrafos))
        atual_linhas += linhas_bloco
    if atual:
        paginas.append(atual)
    return paginas


def agrupar_desligamentos_em_paginas(dados, max_linhas, chars_por_linha):
    """Pagina os parágrafos (com runs) de Desligamentos por orçamento de linhas."""
    paragrafos = gerar_paragrafos_desligamentos(dados)
    paginas, atual, atual_linhas = [], [], 0
    for p in paragrafos:
        linhas = _linhas_ocupadas(p, chars_por_linha)
        if atual and (atual_linhas + linhas > max_linhas):
            paginas.append(atual)
            atual, atual_linhas = [], 0
        atual.append(p)
        atual_linhas += linhas
    if atual:
        paginas.append(atual)
    return paginas or [[]]


def _renderizar_pagina_categoria(prs, titulo, corpo_paragrafos):
    """Cria um slide de conteúdo (clonado do modelo) com título + corpo em texto."""
    novo = _duplicate_slide(prs, 2)
    _remover_fotos(novo)
    shp_categoria = _find_shape(novo, "DESLIGAMENTOS")
    shp_corpo = _find_shape(novo, "Foram registradas")
    if shp_categoria:
        _set_text_preservando_estilo(shp_categoria, titulo)
    if shp_corpo:
        _set_paragrafos_com_estilo(shp_corpo, corpo_paragrafos)
    return novo


def _renderizar_pagina_zeladoria_tabela(prs, titulo, pagina_usinas):
    """
    Cria um slide de Zeladoria com uma TABELA real (Usina | Roçagem | Poda
    Química | Limpeza dos Módulos), no lugar da antiga lista de bullets —
    o texto corrido não cabia numa única página quando havia muitas usinas.
    """
    novo = _duplicate_slide(prs, 2)
    _remover_fotos(novo)
    shp_categoria = _find_shape(novo, "DESLIGAMENTOS")
    shp_corpo = _find_shape(novo, "Foram registradas")
    if shp_categoria:
        _set_text_preservando_estilo(shp_categoria, titulo)

    if shp_corpo:
        fmt = _extrair_formato_base(shp_corpo)
        left, top, width, height = shp_corpo.left, shp_corpo.top, shp_corpo.width, shp_corpo.height
        shp_corpo._element.getparent().remove(shp_corpo._element)
    else:
        fmt = None
        left = top = width = height = None

    if left is None:
        # fallback conservador, caso o placeholder de corpo não seja encontrado
        from pptx.util import Emu
        left, top, width, height = Emu(700000), Emu(1500000), Emu(11500000), Emu(5500000)

    linhas = len(pagina_usinas) + 1  # +1 cabeçalho
    colunas = len(ZELADORIA_COLUNAS)
    graphic_frame = novo.shapes.add_table(linhas, colunas, left, top, width, height)
    tabela = graphic_frame.table

    for j, titulo_col in enumerate(ZELADORIA_COLUNAS):
        cel = tabela.cell(0, j)
        cel.text = titulo_col
        _formatar_celula_tabela(cel, fmt, negrito=True)

    for i, item in enumerate(pagina_usinas, start=1):
        cel_usina = tabela.cell(i, 0)
        cel_usina.text = item["usina"]
        _formatar_celula_tabela(cel_usina, fmt, negrito=True)
        for j, (_, nome_grupo) in enumerate(ZELADORIA_MAPA_COLUNAS, start=1):
            cel = tabela.cell(i, j)
            cel.text = _status_zeladoria_para_tabela(item, nome_grupo)
            _formatar_celula_tabela(cel, fmt, negrito=False)

    return novo


def gerar_relatorio_pptx(cliente, semana_num, data_label, atividades_por_usina,
                          desligamentos_por_usina, usinas_cliente):
    """
    PADRÃO DEFINITIVO confirmado com Fred em 23/07/2026 (revisão meticulosa
    do relatório RENOGRID Semana 30, vale para todos os clientes):

    (1) Capa; (2) Pautas Gerais (5 tópicos fixos); (3-4) ATIVIDADES DA
    SEMANA por usina em ordem alfabética; (5-6) DESLIGAMENTOS, mesma
    lógica; (7) CHAMADOS E PROTOCOLOS (só título — Fred preenche);
    (8) OUTRAS ATIVIDADES (só título); (9) ZELADORIA (estrutura de usinas
    pronta, Fred preenche os valores); (10) contato (slide original do
    template, não gerado).

    atividades_por_usina / desligamentos_por_usina: retorno de
    coletar_atividades_e_desligamentos_por_usina().
    usinas_cliente: lista de usinas do cliente (define a ordem alfabética
    e a estrutura da Zeladoria).
    data_label: mantido por compatibilidade de assinatura (não exibido).
    Retorna BytesIO() pronto para download.
    """
    prs = Presentation(TEMPLATE_PATH)

    # Orçamento de paginação: FIXO, calibrado contra o relatório de
    # referência corrigido pelo Fred (RENOGRID Semana 30 — 5 usinas simples
    # coube confortavelmente numa página). A extração dinâmica a partir do
    # placeholder do modelo (_extrair_orcamento_linhas) foi tentada, mas o
    # placeholder "Foram registradas..." do modelo antigo é baixo (dimensionado
    # para um parágrafo curto + fotos ao lado, removidas por _remover_fotos),
    # o que gerava um orçamento de ~4 linhas — 1 usina por slide. Não usar.
    max_linhas = 16

    usinas_ordenadas = sorted(set(usinas_cliente or []) | set(atividades_por_usina) | set(desligamentos_por_usina),
                               key=_norm)

    # --- Slide 1: Capa — só cliente e semana (sem data) ---------------------
    capa = _duplicate_slide(prs, 0)
    shp_titulo = _find_shape(capa, "O&M")
    if shp_titulo:
        _set_text_preservando_estilo(shp_titulo, f"O&M – {cliente}\nSemana {semana_num}")

    # --- Slide 2: Pautas Gerais — fixo, igual para todos os clientes -------
    _renderizar_pautas_gerais(prs)

    # --- Slides 3-4: ATIVIDADES DA SEMANA (tópico 1) ------------------------
    _renderizar_topico_usinas(prs, 1, "ATIVIDADES DA SEMANA", usinas_ordenadas, atividades_por_usina,
                               _formatar_item_atividade, "Sem atividades realizadas no período.",
                               "•", max_linhas)

    # --- Slides 5-6: DESLIGAMENTOS (tópico 2) -------------------------------
    _renderizar_topico_usinas(prs, 2, "DESLIGAMENTOS", usinas_ordenadas, desligamentos_por_usina,
                               _formatar_item_desligamento, "Sem desligamentos registrados no período.",
                               "•", max_linhas)

    # --- Slide 7: CHAMADOS E PROTOCOLOS — só título, Fred preenche ---------
    def _corpo_chamados(tf):
        _add_paragrafo(tf, "subtitulo", [{"texto": "PROTOCOLOS CONCESSIONÁRIAS", "bold": True}], first=True)
    _renderizar_secao_placeholder(prs, 3, "CHAMADOS E PROTOCOLOS", _corpo_chamados)

    # --- Slide 8: OUTRAS ATIVIDADES — só título, Fred preenche --------------
    _renderizar_secao_placeholder(prs, 4, "OUTRAS ATIVIDADES", None)

    # --- Slide 9: ZELADORIA — estrutura de usinas pronta, Fred preenche ----
    def _corpo_zeladoria(tf):
        primeiro = True
        for i, usina in enumerate(usinas_ordenadas, start=1):
            _add_paragrafo(tf, "usina", [{"texto": usina, "bold": True}], first=primeiro, start_at=i)
            primeiro = False
            _add_paragrafo(tf, "item", [{"texto": "", "bold": False}], first=False, bullet_char="•")
    _renderizar_secao_placeholder(prs, 5, "ZELADORIA", _corpo_zeladoria)

    # --- Reordena o deck: capa nova -> pautas nova -> conteúdo -> contato --
    xml_slides = prs.slides._sldIdLst
    todos_els = list(xml_slides)
    contato_el = todos_els[6]        # slide de contato original do template
    novos_els = todos_els[7:]        # tudo que foi duplicado nesta chamada, já na ordem certa

    for e in todos_els:
        xml_slides.remove(e)
    for e in novos_els + [contato_el]:
        xml_slides.append(e)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
