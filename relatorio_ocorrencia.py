# -*- coding: utf-8 -*-
"""
relatorio_ocorrencia.py
─────────────────────────────────────────────────────────────────────────
Geração do "Relatório de Ocorrência" (.pptx) no modelo Grid Co., a partir
de um resumo curto (expandido pela IA em app.py) + fotos de evidência
enviadas pelo usuário no Painel de Relatórios.

Segue o mesmo padrão estrutural/visual usado no relatório de referência
(Ocorrência UFV Ibaté I, THOPEN) e reaproveita as mesmas convenções de
código do relatorio_semanal.py (_duplicate_slide, _find_shape etc.), mas
como módulo próprio e independente — não altera relatorio_semanal.py.

Slides do modelo (templates/modelo_relatorio_ocorrencia.pptx), 0-based:
  0 = Capa                    ({{CLIENTE}} – {{USINA}})
  1 = Ocorrência               (texto, sem fotos)
  2 = Evidências                (texto + 4 slots de foto lado a lado — duplicado
                                 conforme o número de fotos, 4 por slide)
  3 = Ações a Serem Tomadas    (texto, sem fotos)
  4 = Conclusão                 (texto + assinatura)
  5 = Fale com a gente          (contato, fixo — não é tocado)
"""
import os
import copy
from io import BytesIO

from PIL import Image
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.util import Pt, Emu

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "templates", "modelo_relatorio_ocorrencia.pptx")

# Slots de foto do slide de Evidências (extraídos do modelo original,
# 4 posições lado a lado). Usados em "contain-fit": a foto é redimensionada
# preservando a proporção original e centralizada dentro do slot, nunca
# distorcida nem cortada.
EVIDENCIA_SLOT_Y = 4289026
EVIDENCIA_SLOT_H = 4320000
EVIDENCIA_SLOT_XS = [2716823, 6829333, 10640355, 14451377]
EVIDENCIA_SLOT_W = 3601784
FOTOS_POR_SLIDE = 4

FONT_NAME = "Poppins"
FONT_SIZE = Pt(20)


# ── Helpers reaproveitados do padrão de relatorio_semanal.py ──────────────

def _duplicate_slide(prs, index):
    """Duplica um slide existente do modelo (mesmo layout, cores, fontes, formas)."""
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)

    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    id_map = {}
    for rId, rel in source.part.rels.items():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        new_rid = dest.part.rels._add_relationship(rel.reltype, rel._target, rel.is_external)
        id_map[rId] = new_rid

    for shp in source.shapes:
        newel = copy.deepcopy(shp._element)
        for el in newel.iter():
            for attr in ("embed", "link", "id"):
                full = qn(f"r:{attr}")
                val = el.get(full)
                if val and val in id_map:
                    el.set(full, id_map[val])
        dest.shapes._spTree.append(newel)
    return dest


def _find_shape(slide, contem_texto):
    for shp in slide.shapes:
        if shp.has_text_frame and contem_texto.lower() in shp.text_frame.text.lower():
            return shp
    return None


def _find_shape_by_name(slide, name):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    return None


def _pPr_corpo_xml():
    """Parágrafo de corpo padrão deste template: justificado, 6pt antes,
    sem marcador (igual às seções Ocorrência/Ações/Conclusão do relatório
    de referência)."""
    ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    return parse_xml(
        f'<a:pPr {ns} marL="299720" marR="347980" algn="just">'
        f'<a:spcBef><a:spcPts val="600"/></a:spcBef></a:pPr>'
    )


def _set_paragrafos_corpo(shape, paragrafos):
    """Substitui o conteúdo de um text box por parágrafos de corpo (20pt
    Poppins, justificado), permitindo negrito por trecho dentro da mesma
    linha via 'runs'.

    paragrafos: lista de {"runs": [{"texto": str, "bold": bool}, ...]}
    """
    tf = shape.text_frame
    tf.clear()
    if not paragrafos:
        paragrafos = [{"runs": [{"texto": "", "bold": False}]}]
    for i, par in enumerate(paragrafos):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p_el = p._p
        velho = p_el.find(qn("a:pPr"))
        if velho is not None:
            p_el.remove(velho)
        p_el.insert(0, _pPr_corpo_xml())
        for r in par.get("runs", []):
            run = p.add_run()
            run.text = r.get("texto", "")
            run.font.name = FONT_NAME
            run.font.size = FONT_SIZE
            run.font.bold = bool(r.get("bold", False))


def _remover_pics(slide):
    for shp in list(slide.shapes):
        if shp.shape_type == 13:
            shp._element.getparent().remove(shp._element)


def _set_header(slide, usina):
    shp = _find_shape(slide, "Relatório de Ocorrência")
    if shp:
        for p in shp.text_frame.paragraphs:
            for r in p.runs:
                if "{{USINA}}" in r.text:
                    r.text = r.text.replace("{{USINA}}", usina)


def _add_foto_contida(slide, image_stream, slot_x, slot_y, slot_w, slot_h):
    """Insere uma foto dentro do slot preservando a proporção original
    (contain-fit) e centralizando dentro da área do slot."""
    im = Image.open(image_stream)
    largura_px, altura_px = im.size
    proporcao = largura_px / altura_px
    image_stream.seek(0)

    proporcao_slot = slot_w / slot_h
    if proporcao >= proporcao_slot:
        # foto mais "larga" que o slot -> limita pela largura
        largura = slot_w
        altura = int(slot_w / proporcao)
    else:
        altura = slot_h
        largura = int(slot_h * proporcao)

    left = slot_x + (slot_w - largura) // 2
    top = slot_y + (slot_h - altura) // 2
    slide.shapes.add_picture(image_stream, Emu(int(left)), Emu(int(top)),
                              width=Emu(int(largura)), height=Emu(int(altura)))


# ── Geração principal ──────────────────────────────────────────────────

def gerar_relatorio_ocorrencia_pptx(cliente, usina, ocorrencia_paragrafos,
                                     acoes_paragrafos, conclusao_paragrafos,
                                     fotos_evidencia):
    """
    cliente, usina: strings (ex.: "THOPEN", "Matão II - Topázio").
    ocorrencia_paragrafos / acoes_paragrafos / conclusao_paragrafos:
        lista de {"runs": [{"texto": str, "bold": bool}, ...]} (parágrafos
        de corpo, na ordem em que devem aparecer no slide).
    fotos_evidencia: lista de streams de imagem (BytesIO ou caminho), na
        ordem de upload. Pode ser vazia.
    Retorna BytesIO() pronto para download.
    """
    prs = Presentation(TEMPLATE_PATH)

    # --- Slide 0: Capa -------------------------------------------------
    capa = _duplicate_slide(prs, 0)
    shp_sub = _find_shape(capa, "{{CLIENTE}}")
    if shp_sub:
        for p in shp_sub.text_frame.paragraphs:
            for r in p.runs:
                if "{{CLIENTE}}" in r.text:
                    r.text = r.text.replace("{{CLIENTE}}", cliente).replace("{{USINA}}", usina)

    # --- Slide 1: Ocorrência --------------------------------------------
    s_ocorrencia = _duplicate_slide(prs, 1)
    _set_header(s_ocorrencia, usina)
    shp_corpo = _find_shape_by_name(s_ocorrencia, "CaixaDeTexto 4")
    if shp_corpo:
        _set_paragrafos_corpo(shp_corpo, ocorrencia_paragrafos)

    # --- Slide(s) 2: Evidências (4 fotos por slide) ----------------------
    fotos = fotos_evidencia or []
    total_slides_evidencia = max(1, -(-len(fotos) // FOTOS_POR_SLIDE)) if fotos else 0
    for i in range(total_slides_evidencia):
        s_evid = _duplicate_slide(prs, 2)
        _set_header(s_evid, usina)
        _remover_pics(s_evid)
        shp_corpo_evid = _find_shape_by_name(s_evid, "CaixaDeTexto 1")
        if shp_corpo_evid:
            texto = ([{"runs": [{"texto": "Registro fotográfico da ocorrência:", "bold": False}]}]
                      if i == 0 else [{"runs": [{"texto": "", "bold": False}]}])
            _set_paragrafos_corpo(shp_corpo_evid, texto)
        lote = fotos[i * FOTOS_POR_SLIDE:(i + 1) * FOTOS_POR_SLIDE]
        for slot_idx, foto in enumerate(lote):
            _add_foto_contida(s_evid, foto, EVIDENCIA_SLOT_XS[slot_idx], EVIDENCIA_SLOT_Y,
                               EVIDENCIA_SLOT_W, EVIDENCIA_SLOT_H)

    # --- Slide 3: Ações a Serem Tomadas -----------------------------------
    s_acoes = _duplicate_slide(prs, 3)
    _set_header(s_acoes, usina)
    shp_corpo_acoes = _find_shape_by_name(s_acoes, "CaixaDeTexto 1")
    if shp_corpo_acoes:
        _set_paragrafos_corpo(shp_corpo_acoes, acoes_paragrafos)

    # --- Slide 4: Conclusão ------------------------------------------------
    s_concl = _duplicate_slide(prs, 4)
    _set_header(s_concl, usina)
    shp_corpo_concl = _find_shape_by_name(s_concl, "CaixaDeTexto 4")
    if shp_corpo_concl:
        _set_paragrafos_corpo(shp_corpo_concl, conclusao_paragrafos)

    # --- Reordena: remove os 6 slides originais do modelo, mantém só as
    #     duplicatas geradas nesta chamada + o slide de contato original --
    xml_slides = prs.slides._sldIdLst
    todos_els = list(xml_slides)
    contato_el = todos_els[5]     # slide de contato original do template
    novos_els = todos_els[6:]     # tudo que foi duplicado nesta chamada, já na ordem certa

    for e in todos_els:
        xml_slides.remove(e)
    for e in novos_els + [contato_el]:
        xml_slides.append(e)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
